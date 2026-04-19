import os
import sys
import ast
import io
import json
import re
import shutil
import tempfile
import subprocess
import logging
import time
import hashlib
import requests
import copy
from typing import Dict, List, Optional, Tuple
from concurrent.futures import ThreadPoolExecutor, as_completed
from dotenv import load_dotenv

load_dotenv()

from PyPDF2 import PdfReader
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_PARAGRAPH_ALIGNMENT
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
import asyncio
import edge_tts
from pdf2image import convert_from_path
from google import genai
from PIL import Image, ImageDraw
import colorsys
from docx import Document

from vo_sanitizer import sanitize_voiceover
from diagrams import (
    get_diagram_code, get_scene_code,
    AVAILABLE_DIAGRAM_TYPES, AVAILABLE_SCENE_TYPES, DIAGRAM_SCHEMAS,
)

MANIM_BIN = shutil.which("manim") or os.path.join(sys.prefix, "bin", "manim")

logging.basicConfig(level=logging.INFO, format="%(asctime)s - %(levelname)s - %(message)s")
logger = logging.getLogger(__name__)

OUTPUT_DIR = "outputs"
os.makedirs(OUTPUT_DIR, exist_ok=True)
IMAGE_CACHE_DIR = os.path.join(OUTPUT_DIR, "image_cache")
os.makedirs(IMAGE_CACHE_DIR, exist_ok=True)
PREVIEW_DIR = os.path.join(OUTPUT_DIR, "image_previews")
os.makedirs(PREVIEW_DIR, exist_ok=True)

API_KEY = os.getenv("GEMINI_API_KEY")
if not API_KEY:
    raise ValueError("Set GEMINI_API_KEY environment variable")
GOOGLE_CSE_API_KEY = os.getenv("GOOGLE_CSE_API_KEY")
if not GOOGLE_CSE_API_KEY:
    raise ValueError("Set GOOGLE_CSE_API_KEY environment variable")
GOOGLE_CSE_CX = os.getenv("GOOGLE_CSE_CX")
if not GOOGLE_CSE_CX:
    raise ValueError("Set GOOGLE_CSE_CX environment variable")

client = genai.Client(api_key=API_KEY)

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)
ACCENT_COLORS = [
    RGBColor(0, 217, 255), RGBColor(255, 107, 107),
    RGBColor(78, 205, 196), RGBColor(255, 230, 109),
    RGBColor(167, 139, 250), RGBColor(52, 211, 153),
]
DARK_BG  = RGBColor(15, 23, 42)
TEXT_PRI = RGBColor(241, 245, 249)
TEXT_SEC = RGBColor(148, 163, 184)
TEXT_DIM = RGBColor(71, 85, 105)


# ════════════════════════════════════════════════════════════════════════════
# DOCUMENT LOADING
# ════════════════════════════════════════════════════════════════════════════

def load_pdf_text(path: str) -> str:
    reader = PdfReader(path)
    return "\n".join(page.extract_text() or "" for page in reader.pages)


def load_reference_documents_from_project() -> str:
    texts = []
    excluded = {"venv", "manim_env", "outputs", "__pycache__", ".git", "node_modules"}
    allowed  = {".pdf", ".ppt", ".pptx", ".docx", ".txt", ".md"}
    for root, dirs, files in os.walk(Path(".")):
        dirs[:] = [d for d in dirs if d not in excluded]
        for file in files:
            path = Path(root) / file
            if path.suffix.lower() not in allowed:
                continue
            try:
                ext = path.suffix.lower()
                if ext == ".pdf":
                    texts.append(load_pdf_text(str(path)))
                elif ext in (".ppt", ".pptx"):
                    prs = Presentation(str(path))
                    for slide in prs.slides:
                        for shape in slide.shapes:
                            if shape.has_text_frame:
                                texts.append(shape.text)
                elif ext == ".docx":
                    doc = Document(str(path))
                    texts.extend(p.text for p in doc.paragraphs if p.text.strip())
                elif ext in (".txt", ".md"):
                    with open(path, "r", encoding="utf-8", errors="ignore") as f:
                        texts.append(f.read())
            except Exception as e:
                logger.warning(f"Failed loading {path}: {e}")
    return "\n\n".join(texts)


# ════════════════════════════════════════════════════════════════════════════
# PROMPT (Full AI pipeline)
# ════════════════════════════════════════════════════════════════════════════

def create_prompt(
    topic: str = "",
    transcript_text: str = None,
    slides_text: str = None,
    notes_text: str = None,
) -> str:
    has_t = bool(transcript_text and transcript_text.strip())
    has_s = bool(slides_text and slides_text.strip())
    has_n = bool(notes_text and notes_text.strip())

    t_block = f"""
======================
CLASS TRANSCRIPT (PRIMARY SOURCE)
======================
Follow the EXACT topic order. Preserve key terms and examples.

TRANSCRIPT:
{transcript_text[:20000]}
""" if has_t else ""

    s_block = f"""
======================
SLIDE DECK (STRUCTURE REFERENCE)
======================
SLIDES:
{slides_text[:10000]}
""" if has_s else ""

    n_block = f"""
======================
NOTES / TEXTBOOK (DEPTH REFERENCE)
======================
NOTES:
{notes_text[:8000]}
""" if has_n else ""

    topic_line = f"Topic: {topic}" if topic else "Topic: (infer from uploaded content)"

    return f'''You are an Expert Educational Slide Structure Generator.

{topic_line}

{t_block}{s_block}{n_block}

CRITICAL: ALL content MUST come from the uploaded material. Do NOT invent facts.

TASK: Generate slide structure ONLY — titles, bullets, types, diagrams, image queries.
Voiceovers will be generated in a SEPARATE dedicated pass. Set voiceover to "".

LECTURE MODE:
- Cover ALL major concepts in their original order.
- 20-40 slides total.
- DIAGRAM RULE: type="diagram" only if fundamentally visual AND you have REAL data.
  Max 1 diagram per 4-5 slides. When in doubt -> type="text"

======================
OUTPUT FORMAT — VALID JSON ONLY, NOTHING ELSE
======================

{{
  "metadata": {{
    "total_slides": integer,
    "estimated_duration_minutes": integer,
    "target_audience": "beginner|intermediate|advanced",
    "full_description": "3-4 paragraph academic summary"
  }},
  "slides": [
    {{
      "slide_index": integer,
      "title": "ASCII only, max 55 chars",
      "type": "text|diagram",
      "bullets": ["Substantive key point (max 95 chars, ASCII)", "Supporting point", "Concrete example or fact"],
      "voiceover": "",
      "formula": "LaTeX string if math needed, else empty string",
      "diagram": "diagram type or none",
      "diagram_params": {{}},
      "diagram_needed": boolean,
      "slide_duration_seconds": 90,
      "image": {{
        "needed": boolean,
        "search_queries": ["very specific query 1", "query 2", "query 3"],
        "placement": "right"
      }}
    }}
  ]
}}

AVAILABLE DIAGRAM TYPES:
{', '.join(AVAILABLE_DIAGRAM_TYPES)}

{DIAGRAM_SCHEMAS}

STRICT RULES:
- bullets: 3-5 substantive points per slide, not vague headers.
- type="diagram": fill diagram_params with REAL content from source.
- type="text": set diagram="none", diagram_params={{}}, fill image object.
- image.needed: true only for real-world objects/experimental setups.
- image.search_queries: VERY specific (e.g. "ResNet50 architecture diagram").
- voiceover: always set to "" — do not generate it here.
- All strings: ASCII only, no unicode.
- Output ONLY the JSON. No markdown, no explanation.
'''


# ════════════════════════════════════════════════════════════════════════════
# GEMINI VISION RELEVANCE CHECK
# ════════════════════════════════════════════════════════════════════════════

def gemini_vision_check(image_path: str, slide_title: str, topic: str = "") -> bool:
    try:
        img = Image.open(image_path).convert("RGB")
        img.thumbnail((512, 512), Image.LANCZOS)
        context = f"a lecture slide titled '{slide_title}'"
        if topic:
            context += f" in a course about '{topic}'"
        prompt = (
            f"This image is being placed on {context}. "
            f"Does this image directly and visually relate to that content? "
            f"Reply with only YES or NO."
        )
        response = client.models.generate_content(
            model="gemini-2.5-flash",
            contents=[prompt, img],
        )
        result = response.text.strip().upper()
        is_relevant = result.startswith("YES")
        logger.info(
            f"  Vision '{slide_title[:28]}': "
            f"{'✓ RELEVANT' if is_relevant else '✗ REJECTED'} ({result[:10]})"
        )
        return is_relevant
    except Exception as e:
        logger.warning(f"  Vision check failed (accepting by default): {e}")
        return True


# ════════════════════════════════════════════════════════════════════════════
# IMAGE SEARCH + CACHE
# ════════════════════════════════════════════════════════════════════════════

def _cache_key(query: str) -> str:
    return os.path.join(IMAGE_CACHE_DIR, hashlib.md5(query.encode()).hexdigest() + ".jpg")


def invalidate_cse_cache_for_queries(queries: List[str]) -> None:
    for q in queries:
        cached = _cache_key(q)
        if os.path.exists(cached):
            try:
                os.remove(cached)
                logger.info(f"  CSE cache cleared: {q[:45]}")
            except Exception as e:
                logger.warning(f"  Could not clear CSE cache '{q[:30]}': {e}")


def search_google_cse_images(query: str, bypass_cache: bool = False) -> Optional[str]:
    cached = _cache_key(query)
    if not bypass_cache and os.path.exists(cached) and os.path.getsize(cached) > 2000:
        logger.info(f"  Cache hit: {query[:45]}")
        return cached

    logger.info(f"  CSE {'(fresh)' if bypass_cache else ''}: {query[:55]}")
    try:
        params = {
            "key": GOOGLE_CSE_API_KEY,
            "cx": GOOGLE_CSE_CX,
            "q": query,
            "searchType": "image",
            "num": 8,
            "safe": "active",
            "imgType": "photo",
            "imgSize": "large",
        }
        if bypass_cache:
            params["start"] = 2

        r = requests.get(
            "https://www.googleapis.com/customsearch/v1",
            params=params,
            timeout=12,
        )
        r.raise_for_status()
        data = r.json()
        if "error" in data:
            logger.warning(f"  CSE error: {data['error'].get('message','')}")
            return None
        for item in data.get("items", []):
            url = item.get("link", "")
            if not url or any(url.lower().endswith(e) for e in [".svg", ".gif", ".webp"]):
                continue
            if _download_image(url, cached):
                return cached
            time.sleep(0.15)
    except Exception as e:
        logger.warning(f"  CSE failed: {e}")
    return None


def _download_image(url: str, save_path: str, retries: int = 2) -> bool:
    for _ in range(retries):
        try:
            r = requests.get(url, timeout=12,
                             headers={"User-Agent": "Mozilla/5.0", "Accept": "image/*"},
                             stream=True)
            r.raise_for_status()
            ct = r.headers.get("Content-Type", "")
            if any(x in ct for x in ["svg", "html", "text"]):
                return False
            content = r.content
            if len(content) < 5000:
                return False
            with open(save_path, "wb") as f:
                f.write(content)
            img = Image.open(save_path)
            if img.width < 200 or img.height < 200:
                os.remove(save_path); return False
            img.convert("RGB").save(save_path, "JPEG", quality=88)
            return True
        except Exception:
            if os.path.exists(save_path):
                os.remove(save_path)
            time.sleep(0.3)
    return False


def _procedural_image(work_dir: str, slide_idx: int, title: str, keywords: List[str]) -> Optional[str]:
    try:
        w, h = 800, 600
        hue = abs(hash(" ".join(keywords))) % 360 / 360.0
        rgb = colorsys.hsv_to_rgb(hue, 0.5, 0.3)
        bg  = tuple(int(c * 255) for c in rgb)
        img = Image.new("RGB", (w, h), color=bg)
        draw = ImageDraw.Draw(img)
        for y in range(h):
            c = tuple(min(255, int(bg[i] * (0.8 + 0.2 * y / h))) for i in range(3))
            draw.line([(0, y), (w, y)], fill=c)
        draw.text((w // 2, 60), title[:45], fill=(200, 220, 255), anchor="mm")
        for i, kw in enumerate(keywords[:4]):
            cx, cy, r = 100 + i * 170, h // 2, 55
            draw.ellipse([cx-r, cy-r, cx+r, cy+r], outline=(170, 190, 220), width=2)
            draw.text((cx, cy + r + 14), kw[:12], fill=(160, 180, 210), anchor="mm")
        path = os.path.join(work_dir, f"proc_{slide_idx}.jpg")
        img.save(path, "JPEG", quality=82)
        return path
    except Exception as e:
        logger.warning(f"Procedural image failed: {e}")
        return None


def _find_image(
    slide_idx: int,
    image_spec: Dict,
    work_dir: str,
    slide_title: str = "",
    topic: str = "",
    bypass_cache: bool = False,
) -> Optional[str]:
    if not image_spec or not image_spec.get("needed"):
        return None
    queries = [q for q in image_spec.get("search_queries", []) if q.strip()]
    if not queries:
        return None
    label = slide_title or queries[0]
    for i, query in enumerate(queries):
        cached = search_google_cse_images(query, bypass_cache=bypass_cache)
        if cached:
            if gemini_vision_check(cached, label, topic):
                dest = os.path.join(work_dir, f"img_{slide_idx}.jpg")
                shutil.copy2(cached, dest)
                logger.info(f"Slide {slide_idx+1}: image ready (q{i+1}, vision ✓)")
                return dest
            else:
                logger.info(f"Slide {slide_idx+1}: q{i+1} rejected, trying next")
        time.sleep(0.4)
    logger.warning(f"Slide {slide_idx+1}: all queries rejected → procedural fallback")
    return _procedural_image(work_dir, slide_idx, label, queries)


def prefetch_images(
    slides: List[Dict],
    work_dir: str,
    max_workers: int = 4,
    topic: str = "",
    skip_indices: set = None,
) -> Dict[int, str]:
    skip_indices = skip_indices or set()
    tasks = [
        (i, slide.get("image", {}), slide.get("title", ""))
        for i, slide in enumerate(slides)
        if slide.get("type") == "text"
        and slide.get("image", {}).get("needed")
        and i not in skip_indices
    ]
    if not tasks:
        return {}
    logger.info(f"Fetching {len(tasks)} images in parallel (Vision filtering)...")

    def fetch(args):
        idx, image_spec, slide_title = args
        return idx, _find_image(idx, image_spec, work_dir,
                                slide_title=slide_title, topic=topic)

    results = {}
    with ThreadPoolExecutor(max_workers=max_workers) as ex:
        for idx, path in [f.result() for f in as_completed(
            {ex.submit(fetch, t): t for t in tasks}
        )]:
            if path:
                results[idx] = path
    logger.info(f"Images ready: {len(results)}/{len(tasks)}")
    return results


# ════════════════════════════════════════════════════════════════════════════
# PREVIEW IMAGE
# ════════════════════════════════════════════════════════════════════════════

def fetch_preview_image(
    slide_idx: int,
    slide: Dict,
    topic: str = "",
    force_refresh: bool = False,
) -> Optional[str]:
    preview_path = os.path.join(PREVIEW_DIR, f"preview_{slide_idx}.jpg")
    image_spec   = slide.get("image", {})
    queries      = [q for q in image_spec.get("search_queries", []) if q.strip()]
    slide_title  = slide.get("title", queries[0] if queries else f"Slide {slide_idx + 1}")

    if force_refresh:
        if os.path.exists(preview_path):
            try:
                os.remove(preview_path)
            except Exception:
                pass
        invalidate_cse_cache_for_queries(queries)
        logger.info(f"Force-refresh slide {slide_idx + 1}: all caches cleared")
    else:
        if os.path.exists(preview_path) and os.path.getsize(preview_path) > 2000:
            logger.info(f"Preview cache hit for slide {slide_idx + 1}")
            return preview_path

    if not queries:
        logger.info(f"Slide {slide_idx + 1}: no image queries defined")
        return None

    for i, query in enumerate(queries):
        cached = search_google_cse_images(query, bypass_cache=force_refresh)
        if cached:
            if gemini_vision_check(cached, slide_title, topic):
                shutil.copy2(cached, preview_path)
                logger.info(f"Preview ready slide {slide_idx + 1} (q{i + 1}, vision ✓)")
                return preview_path
            else:
                logger.info(f"Preview q{i + 1} rejected slide {slide_idx + 1}, trying next")
        time.sleep(0.3)

    logger.warning(f"No relevant preview found for slide {slide_idx + 1}")
    return None


def invalidate_preview_cache(slide_idx: int) -> None:
    preview_path = os.path.join(PREVIEW_DIR, f"preview_{slide_idx}.jpg")
    if os.path.exists(preview_path):
        try:
            os.remove(preview_path)
            logger.info(f"Preview cache cleared for slide {slide_idx + 1}")
        except Exception as e:
            logger.warning(f"Could not clear preview cache for slide {slide_idx + 1}: {e}")


def _fallback_image(work_dir: str, slide_idx: int, title: str) -> str:
    img  = Image.new("RGB", (1920, 1080), color=(15, 23, 42))
    draw = ImageDraw.Draw(img)
    draw.text((960, 440), title[:75], fill=(0, 217, 255), anchor="mm")
    draw.text((960, 540), "[Diagram rendering failed]", fill=(148, 163, 184), anchor="mm")
    path = os.path.join(work_dir, f"fallback_{slide_idx}.png")
    img.save(path)
    return path


# ════════════════════════════════════════════════════════════════════════════
# PPTX CREATION  (Full AI pipeline)
# ════════════════════════════════════════════════════════════════════════════

def _tb(slide_obj, left, top, width, height, text, font_size, color,
        bold=False, wrap=True, align=PP_PARAGRAPH_ALIGNMENT.LEFT):
    tb = slide_obj.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.text = text; p.alignment = align
    p.font.size = font_size; p.font.bold = bold
    p.font.color.rgb = color; p.font.name = "Calibri"
    return tb


def create_pptx_with_animations(slides: List[Dict], work_dir: str,
                                  prefetched_images: Dict = None) -> str:
    logger.info("Building PowerPoint...")
    prefetched_images = prefetched_images or {}
    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H

    for i, slide in enumerate(slides):
        try:
            so = prs.slides.add_slide(prs.slide_layouts[6])
            bg = so.background.fill
            bg.solid(); bg.fore_color.rgb = DARK_BG
            accent = ACCENT_COLORS[i % len(ACCENT_COLORS)]

            bar = so.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), SLIDE_W, Inches(0.07))
            bar.fill.solid(); bar.fill.fore_color.rgb = accent; bar.line.fill.background()

            _tb(so, Inches(11.9), Inches(0.1), Inches(1.3), Inches(0.28),
                f"{i+1}/{len(slides)}", Pt(9), TEXT_DIM, align=PP_PARAGRAPH_ALIGNMENT.RIGHT)

            _tb(so, Inches(0.45), Inches(0.17), Inches(12.6), Inches(0.72),
                slide.get("title", f"Slide {i+1}")[:55], Pt(30), accent, bold=True)

            sep = so.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                       Inches(0.45), Inches(0.95), Inches(12.43), Inches(0.018))
            sep.fill.solid(); sep.fill.fore_color.rgb = RGBColor(30, 41, 59); sep.line.fill.background()

            if slide.get("type") == "diagram":
                dtype = slide.get("diagram", "diagram").replace("_", " ").title()
                _tb(so, Inches(2.0), Inches(3.0), Inches(9.0), Inches(1.2),
                    f"[ {dtype} Animation ]", Pt(22), TEXT_SEC,
                    align=PP_PARAGRAPH_ALIGNMENT.CENTER)
            else:
                image_path = prefetched_images.get(i)
                has_img    = image_path and os.path.exists(image_path)

                if has_img:
                    bullet_left  = Inches(0.45)
                    bullet_width = Inches(7.1)
                    img_left, img_top = Inches(7.85), Inches(1.1)
                    img_zone_w, img_zone_h = Inches(5.0), Inches(6.1)
                    try:
                        with Image.open(image_path) as im:
                            iw, ih = im.size
                        aspect = iw / ih
                        if aspect > 5.0 / 6.1:
                            fw = img_zone_w
                            fh = Emu(int(img_zone_w / aspect))
                        else:
                            fh = img_zone_h
                            fw = Emu(int(img_zone_h * aspect))
                        vo = (img_zone_h - fh) / 2
                        ho = (img_zone_w - fw) / 2
                        so.shapes.add_picture(image_path, img_left + ho, img_top + vo,
                                              width=fw, height=fh)
                    except Exception as e:
                        logger.warning(f"  Slide {i+1}: image placement failed — {e}")
                        has_img = False
                        bullet_width = Inches(12.43)
                else:
                    bullet_left  = Inches(0.45)
                    bullet_width = Inches(12.43)

                bullets = slide.get("bullets", [])[:5]
                if not bullets and not has_img:
                    # Title/Intro/Outro slides often have no bullets. Fallback to voiceover text
                    # or title if voiceover is also missing. Break into sentences to simulate bullets.
                    vo_text = slide.get("voiceover", "").strip() or slide.get("title", "")
                    # Split by common sentence endings and filter empties
                    potential_bullets = [s.strip() for s in re.split(r'(?<=[.!?])\s+', vo_text)]
                    bullets = [b for b in potential_bullets if b][:5]
                
                y0, row_h = 1.15, 0.92
                for j, bullet in enumerate(bullets):
                    bullet = str(bullet).strip()
                    if len(bullet) > 105: bullet = bullet[:102] + "..."
                    y = Inches(y0 + j * row_h)
                    dot = so.shapes.add_shape(MSO_SHAPE.OVAL,
                                              bullet_left, y + Inches(0.3), Inches(0.12), Inches(0.12))
                    dot.fill.solid(); dot.fill.fore_color.rgb = accent; dot.line.fill.background()
                    _tb(so, bullet_left + Inches(0.23), y,
                        bullet_width - Inches(0.3), Inches(row_h - 0.05),
                        bullet, Pt(19), TEXT_PRI, wrap=True)

            fb = so.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(7.3), SLIDE_W, Inches(0.2))
            fb.fill.solid(); fb.fill.fore_color.rgb = RGBColor(10, 16, 30); fb.line.fill.background()
            _tb(so, Inches(0.45), Inches(7.3), Inches(6), Inches(0.18),
                "Educational Video", Pt(8), TEXT_DIM)

        except Exception as e:
            logger.warning(f"Slide {i+1} PPTX error: {e}")

    path = os.path.join(work_dir, "presentation.pptx")
    prs.save(path)
    logger.info(f"PPTX saved ({len(slides)} slides)")
    return path


def pptx_to_images(pptx_path: str, work_dir: str) -> Dict[int, str]:
    logger.info("Converting PPTX -> images (LibreOffice)...")
    try:
        soffice = shutil.which("soffice") or "/Applications/LibreOffice.app/Contents/MacOS/soffice"
        if not soffice or not os.path.exists(soffice):
            raise FileNotFoundError("LibreOffice not found")

        profile_dir = os.path.join(work_dir, "lo_profile")
        os.makedirs(profile_dir, exist_ok=True)
        user_profile_uri = Path(profile_dir).resolve().as_uri()

        pdf_path = os.path.join(
            work_dir,
            os.path.splitext(os.path.basename(pptx_path))[0] + ".pdf"
        )

        result = subprocess.run(
            [
                soffice,
                "--headless", "--nologo", "--nodefault", "--norestore", "--nolockcheck",
                f"-env:UserInstallation={user_profile_uri}",
                "--convert-to", "pdf:impress_pdf_Export",
                pptx_path, "--outdir", work_dir,
            ],
            check=True, capture_output=True, timeout=180,
        )
        if not os.path.exists(pdf_path):
            logger.error(f"LibreOffice stderr: {result.stderr.decode()[:400]}")
            raise FileNotFoundError("PDF not created")
        image_paths = {}
        for idx, img in enumerate(convert_from_path(pdf_path, dpi=150)):
            p = os.path.join(work_dir, f"slide_{idx}.png")
            img.save(p, "PNG")
            image_paths[idx] = p
        logger.info(f"Converted {len(image_paths)} slides")
        return image_paths
    except Exception as e:
        logger.error(f"PPTX->images failed: {e}")
        return {}


# ════════════════════════════════════════════════════════════════════════════
# MANIM RENDERING (shared by all pipelines)
# ════════════════════════════════════════════════════════════════════════════

def _try_render(code: str, slide_idx: int, work_dir: str, env: dict) -> Tuple[Optional[str], Optional[str]]:
    scene_file = os.path.join(work_dir, f"scene_{slide_idx}.py")
    with open(scene_file, "w", encoding="utf-8") as f:
        f.write(code)
    result = subprocess.run(
        [MANIM_BIN, "-qm", "--disable_caching", "--frame_rate", "30",
         scene_file, f"Slide{slide_idx}"],
        cwd=work_dir, capture_output=True, text=True, timeout=300, env=env,
    )
    if result.returncode != 0:
        return None, (result.stderr or result.stdout)[-1500:]
    media_dir = os.path.join(work_dir, "media", "videos")
    if os.path.exists(media_dir):
        for root, _, files in os.walk(media_dir):
            for fn in files:
                if fn == f"Slide{slide_idx}.mp4":
                    return os.path.join(root, fn), None
    return None, "MP4 not found after render"


def _validate_code(code: str, slide_idx: int) -> Tuple[bool, str]:
    try:
        ast.parse(code)
    except SyntaxError as e:
        return False, str(e)
    except SystemError:
        # Python 3.11 bug: AST recursion depth mismatch on deeply nested code.
        # The code is likely valid — let Manim attempt the render.
        logger.warning(f"  Slide {slide_idx}: ast.parse SystemError (Python bug), skipping AST check")
    if f"class Slide{slide_idx}" not in code:
        return False, f"Missing class Slide{slide_idx}"
    if "def construct(self)" not in code:
        return False, "Missing construct(self)"
    return True, "ok"


def render_diagrams(slides: List[Dict], work_dir: str) -> Dict[int, str]:
    logger.info("Rendering Manim diagrams...")
    if not MANIM_BIN or not os.path.exists(MANIM_BIN):
        logger.warning(f"Manim not found at '{MANIM_BIN}'")
        return {}
    env = os.environ.copy()
    env["PATH"] = "/Library/TeX/texbin:" + env.get("PATH", "")
    try:
        r = subprocess.run([MANIM_BIN, "--version"],
                           capture_output=True, text=True, timeout=120, env=env)
        if r.returncode != 0:
            logger.warning(f"Manim version check failed: {r.stderr[:200]}")
            return {}
        logger.info(f"Manim ready: {r.stdout.strip()}")
    except Exception as e:
        logger.warning(f"Manim unavailable: {e}")
        return {}

    diagram_slides = [
        (i, s) for i, s in enumerate(slides)
        if s.get("type") == "diagram" and s.get("diagram", "none") not in ("none", "", None)
    ]
    if not diagram_slides:
        return {}

    logger.info(f"Rendering {len(diagram_slides)} diagrams...")
    outputs: Dict[int, str] = {}

    def render_one(args):
        i, slide = args
        dtype    = slide.get("diagram", "block_flow")
        duration = int(slide.get("slide_duration_seconds", 90))
        params   = dict(slide.get("diagram_params") or {})
        params["type"]  = dtype
        params["title"] = params.get("title") or slide.get("title", "")
        slide_dir = os.path.join(work_dir, f"manim_{i}")
        os.makedirs(slide_dir, exist_ok=True)
        code = get_diagram_code(i, params, duration)
        if not code:
            return i, None
        ok, err = _validate_code(code, i)
        if not ok:
            logger.error(f"  Slide {i+1}: invalid code: {err}")
            return i, None
        video, render_err = _try_render(code, i, slide_dir, env)
        if not video:
            logger.warning(f"  Slide {i+1}: failed — {(render_err or '')[:120]}")
        return i, video

    with ThreadPoolExecutor(max_workers=2) as ex:
        futures = {ex.submit(render_one, arg): arg[0] for arg in diagram_slides}
        for future in as_completed(futures):
            i, path = future.result()
            if path:
                outputs[i] = path

    logger.info(f"Diagrams done: {len(outputs)}/{len(diagram_slides)}")
    return outputs


# ════════════════════════════════════════════════════════════════════════════
# VOICEOVERS  (Edge TTS — used by all pipelines)
# ════════════════════════════════════════════════════════════════════════════

EDGE_TTS_VOICE = "en-IN-NeerjaNeural"


def _gen_one_vo(args) -> Tuple[int, Dict]:
    idx, slide, work_dir = args
    try:
        time.sleep(idx * 0.07)

        # Get raw voiceover text
        text = slide.get("voiceover", "").strip()
        if not text:
            text = " ".join(slide.get("bullets", [])) or slide.get("title", f"Slide {idx+1}")

        # ── Sanitize: strip markdown that TTS reads literally ──────
        text = sanitize_voiceover(text)
        # ──────────────────────────────────────────────────────────

        path = os.path.join(work_dir, f"audio_{idx}.mp3")

        async def _synth():
            communicate = edge_tts.Communicate(text, EDGE_TTS_VOICE)
            await communicate.save(path)

        asyncio.run(_synth())
        dur = _duration(path)
        logger.info(f"Slide {idx+1}: audio {dur:.1f}s  [{EDGE_TTS_VOICE}]")
        return idx, {"path": path, "duration": dur}
    except Exception as e:
        logger.error(f"Slide {idx+1} audio failed: {e}")
        return idx, {"path": "", "duration": 8.0}


def generate_voiceovers(slides: List[Dict], work_dir: str, max_workers: int = 4) -> Dict[int, Dict]:
    logger.info(f"Generating voiceovers in parallel (Edge TTS · {EDGE_TTS_VOICE})...")
    args = [(i, s, work_dir) for i, s in enumerate(slides)]
    audio_map: Dict[int, Dict] = {}
    with ThreadPoolExecutor(max_workers=max_workers) as ex:
        for idx, result in [f.result() for f in as_completed(
            {ex.submit(_gen_one_vo, a): a for a in args}
        )]:
            audio_map[idx] = result
    logger.info(f"Voiceovers done: {len(audio_map)}")
    return audio_map


def _duration(path: str) -> float:
    try:
        r = subprocess.run(
            ["ffprobe", "-v", "error", "-show_entries", "format=duration",
             "-of", "default=noprint_wrappers=1:nokey=1", path],
            capture_output=True, text=True, timeout=10,
        )
        if r.stdout.strip():
            return float(r.stdout.strip())
    except Exception:
        pass
    return 8.0


# ════════════════════════════════════════════════════════════════════════════
# VIDEO COMPOSITION  (Full AI pipeline)
# ════════════════════════════════════════════════════════════════════════════

def compose_hybrid_video(
    slides: List[Dict],
    image_paths: Dict[int, str],
    audio_map: Dict[int, Dict],
    diagrams: Dict[int, str],
    work_dir: str,
) -> str:
    logger.info("Composing final video...")
    clip_files = []
    TARGET_FPS = 30

    for i, slide in enumerate(slides):
        try:
            info    = audio_map.get(i, {})
            aud     = info.get("path", "")
            aud_ok  = bool(aud and os.path.exists(aud))
            duration = float(
                info.get("duration", slide.get("slide_duration_seconds", 8))
                if aud_ok else slide.get("slide_duration_seconds", 8)
            )

            is_video = False
            if slide.get("type") == "diagram" and i in diagrams and os.path.exists(diagrams[i]):
                source = diagrams[i]; is_video = True
            elif i in image_paths and os.path.exists(image_paths[i]):
                source = image_paths[i]
            else:
                source = _fallback_image(work_dir, i, slide.get("title", f"Slide {i+1}"))

            if not source or not os.path.exists(source):
                logger.warning(f"Slide {i+1}: no visual source")
                continue

            intermediate = os.path.join(work_dir, f"clip_{i}.mp4")
            cmd = ["ffmpeg", "-y", "-loglevel", "error"]
            if is_video:
                cmd += ["-stream_loop", "-1", "-i", source]
            else:
                cmd += ["-loop", "1", "-i", source]
            if aud_ok:
                cmd += ["-i", aud]
            else:
                cmd += ["-f", "lavfi", "-i", "anullsrc=r=44100:cl=stereo"]
            cmd += [
                "-vf",
                "scale=1920:1080:force_original_aspect_ratio=decrease,"
                "pad=1920:1080:(ow-iw)/2:(oh-ih)/2:color=#0F172A",
                "-r", str(TARGET_FPS),
                "-c:v", "libx264", "-preset", "medium", "-crf", "20", "-pix_fmt", "yuv420p",
                "-c:a", "aac", "-b:a", "192k", "-ar", "44100", "-ac", "2",
                "-t", f"{duration:.2f}", "-shortest",
                intermediate,
            ]
            subprocess.run(cmd, check=True, capture_output=True, timeout=300)

            if os.path.exists(intermediate) and os.path.getsize(intermediate) > 5000:
                clip_files.append(intermediate)
                logger.info(f"Clip {i+1}: {duration:.1f}s @{TARGET_FPS}fps")
            else:
                logger.warning(f"Clip {i+1}: output empty/missing")

        except Exception as e:
            logger.error(f"Slide {i+1}: {e}")

    if not clip_files:
        raise ValueError("No clips created!")

    concat_file = os.path.join(work_dir, "concat.txt")
    with open(concat_file, "w") as f:
        for clip in clip_files:
            f.write(f"file '{clip}'\n")

    output = os.path.join(work_dir, "final.mp4")
    subprocess.run(
        ["ffmpeg", "-y", "-f", "concat", "-safe", "0", "-i", concat_file, "-c", "copy", output],
        check=True, capture_output=True, timeout=600,
    )
    if not os.path.exists(output):
        raise FileNotFoundError("Final video not created")

    logger.info(f"Final video: {os.path.getsize(output)/(1024*1024):.1f} MB")
    return output


# ════════════════════════════════════════════════════════════════════════════
# STRUCTURE GENERATION  (Full AI pipeline)
# ════════════════════════════════════════════════════════════════════════════

def _repair_json(text: str) -> str:
    text = re.sub(r"^```(?:json)?", "", text, flags=re.MULTILINE)
    text = re.sub(r"```\s*$", "", text, flags=re.MULTILINE)
    text = text.strip()
    start = text.find("{")
    end   = text.rfind("}")
    if start != -1 and end != -1 and end > start:
        text = text[start:end + 1]
    text = re.sub(r"(?<!\\)\\(?![\\\"nrtbf/u0-9])", r"\\\\", text)
    text = re.sub(r",\s*([\]\}])", r"\1", text)
    text = re.sub(r"\bTrue\b",  "true",  text)
    text = re.sub(r"\bFalse\b", "false", text)
    text = re.sub(r"\bNone\b",  "null",  text)
    return text


def generate_lecture_structure(
    topic: str = "",
    transcript_text: str = None,
    slides_text: str = None,
    notes_text: str = None,
) -> Dict:
    logger.info(f"Generating structure (pass 1): {topic or '(from content)'}")
    prompt = create_prompt(topic, transcript_text, slides_text, notes_text)

    structure = None
    for attempt in range(3):
        response = client.models.generate_content(model="gemini-2.5-flash", contents=prompt)
        raw = response.text.strip()
        text = _repair_json(raw)
        try:
            structure = json.loads(text)
            logger.info(
                f"Structure parsed OK (attempt {attempt + 1}): "
                f"{structure['metadata']['total_slides']} slides, "
                f"~{structure['metadata']['estimated_duration_minutes']} min"
            )
            break
        except json.JSONDecodeError as e:
            logger.warning(f"JSON parse attempt {attempt + 1}/3 failed: {e}")
            if attempt == 2:
                err_pos = getattr(e, "pos", 0)
                snippet = raw[max(0, err_pos - 80): err_pos + 80]
                logger.error(f"Failed JSON snippet:\n{snippet}")
                raise RuntimeError(
                    f"Gemini returned unparseable JSON after 3 attempts. "
                    f"Error at char {err_pos}: {e.msg}\nSnippet: ...{snippet}..."
                ) from e
            time.sleep(2)

    logger.info("Generating voiceovers (pass 2 — dedicated narration pass)…")
    structure["slides"] = _enrich_voiceovers(
        slides          = structure["slides"],
        transcript_text = transcript_text,
        notes_text      = notes_text,
        context         = None,
        topic           = topic,
    )
    return structure


# ════════════════════════════════════════════════════════════════════════════
# SINGLE SLIDE REGENERATION  (Full AI pipeline)
# ════════════════════════════════════════════════════════════════════════════

def create_single_slide_prompt(
    slide_idx, total_slides, current_slide, surrounding_slides,
    transcript_text, slides_text, notes_text, mode, custom_prompt,
) -> str:
    ctx_parts = []
    if transcript_text and transcript_text.strip():
        ctx_parts.append(f"CLASS TRANSCRIPT:\n{transcript_text[:5000]}")
    if slides_text and slides_text.strip():
        ctx_parts.append(f"SLIDE DECK:\n{slides_text[:3000]}")
    if notes_text and notes_text.strip():
        ctx_parts.append(f"NOTES:\n{notes_text[:3000]}")
    ctx = "\n\n---\n\n".join(ctx_parts) if ctx_parts else "No reference material provided."

    surrounding_info = "\n".join(
        f"  Slide {s.get('slide_index','?')}: \"{s.get('title','(untitled)')}\" [{s.get('type','text')}]"
        for s in surrounding_slides
    ) or "  (none)"

    mode_instruction = (
        "MODE: VOICEOVER ONLY\nRewrite ONLY the 'voiceover' field. Keep ALL other fields EXACTLY as-is."
        if mode == "voiceover_only"
        else "MODE: FULL REGENERATION\nRegenerate title, bullets, voiceover, type. Keep slide_index unchanged."
    )
    custom_section = (
        f"\n{'='*50}\nSPECIAL INSTRUCTIONS (highest priority):\n{custom_prompt.strip()}\n{'='*50}\n"
        if custom_prompt and custom_prompt.strip() else ""
    )

    return f"""You are regenerating ONE slide in an educational lecture video.

POSITION: Slide {slide_idx + 1} of {total_slides}
SURROUNDING: {surrounding_info}
CURRENT SLIDE: {json.dumps(current_slide, indent=2)}
SOURCE: {ctx}

{mode_instruction}
{custom_section}
Return ONLY a JSON object with fields: slide_index, title, type, bullets, voiceover,
formula, diagram, diagram_params, diagram_needed, slide_duration_seconds, image.
No markdown, no explanation.
"""


def regenerate_single_slide(
    slide_idx: int,
    current_slide: Dict,
    structure_metadata: Dict,
    surrounding_slides: List[Dict],
    transcript_text: str = None,
    slides_text: str = None,
    notes_text: str = None,
    mode: str = "full",
    custom_prompt: str = "",
) -> Dict:
    logger.info(f"Regenerating slide {slide_idx + 1} (mode={mode})")
    total  = structure_metadata.get("total_slides", 20)
    prompt = create_single_slide_prompt(
        slide_idx, total, current_slide, surrounding_slides,
        transcript_text or "", slides_text or "", notes_text or "",
        mode, custom_prompt,
    )
    response = client.models.generate_content(model="gemini-2.5-flash", contents=prompt)
    text = response.text.strip()
    text = re.sub(r"^```(?:json)?", "", text, flags=re.MULTILINE)
    text = re.sub(r"```\s*$", "", text, flags=re.MULTILINE)
    text = text.strip()
    try:
        result = json.loads(text)
    except json.JSONDecodeError:
        text = re.sub(r"(?<!\\)\\(?![\\\"nrtbf/u])", r"\\\\", text)
        result = json.loads(text)

    if mode == "voiceover_only":
        merged = copy.deepcopy(current_slide)
        merged["voiceover"] = result.get("voiceover", current_slide.get("voiceover", ""))
        result = merged

    result["slide_index"] = current_slide.get("slide_index", slide_idx)
    logger.info(f"Slide {slide_idx + 1} regenerated OK")
    return result


# ════════════════════════════════════════════════════════════════════════════
# FULL AI PIPELINE ORCHESTRATION
# ════════════════════════════════════════════════════════════════════════════

def create_lecture(
    topic: str = "",
    transcript_text: str = None,
    slides_text: str = None,
    notes_text: str = None,
    skip_diagrams: bool = False,
    slides_override: List[Dict] = None,
    custom_images_bytes: Dict[int, bytes] = None,
) -> str:
    logger.info("=" * 70)
    logger.info(f"Creating lecture: {topic or '(from content)'}")
    logger.info("=" * 70)

    with tempfile.TemporaryDirectory() as work_dir:
        if slides_override is not None:
            logger.info(f"Using {len(slides_override)} user-reviewed slides")
            slides = slides_override
            est_dur = max(1, sum(s.get("slide_duration_seconds", 90) for s in slides) // 60)
            meta = {
                "total_slides": len(slides),
                "estimated_duration_minutes": est_dur,
                "target_audience": "intermediate",
                "full_description": (
                    f"Lecture on {topic or 'the uploaded content'} — "
                    f"{len(slides)} slides, approximately {est_dur} minutes."
                ),
            }
            structure = {"metadata": meta, "slides": slides}
        else:
            structure = generate_lecture_structure(topic, transcript_text, slides_text, notes_text)
            meta   = structure["metadata"]
            slides = structure["slides"]

        logger.info(f"Slides: {meta['total_slides']} | ~{meta['estimated_duration_minutes']} min")

        custom_paths: Dict[int, str] = {}
        if custom_images_bytes:
            for idx, img_bytes in custom_images_bytes.items():
                if not img_bytes:
                    continue
                try:
                    cpath = os.path.join(work_dir, f"custom_{idx}.jpg")
                    img = Image.open(io.BytesIO(img_bytes)).convert("RGB")
                    img.save(cpath, "JPEG", quality=90)
                    custom_paths[idx] = cpath
                    logger.info(f"Custom image saved for slide {idx + 1}")
                except Exception as e:
                    logger.warning(f"Failed to save custom image for slide {idx + 1}: {e}")

        prefetched = prefetch_images(
            slides, work_dir, topic=topic,
            skip_indices=set(custom_paths.keys()),
        )
        prefetched.update(custom_paths)

        pptx_path   = create_pptx_with_animations(slides, work_dir, prefetched_images=prefetched)
        image_paths = pptx_to_images(pptx_path, work_dir)

        diagrams:  Dict[int, str] = {}
        audio_map: Dict[int, Dict] = {}

        if skip_diagrams:
            logger.info("Skipping Manim")
            audio_map = generate_voiceovers(slides, work_dir)
        else:
            with ThreadPoolExecutor(max_workers=2) as outer:
                f_diag  = outer.submit(render_diagrams, slides, work_dir)
                f_audio = outer.submit(generate_voiceovers, slides, work_dir)
                diagrams  = f_diag.result()
                audio_map = f_audio.result()

        video_path = compose_hybrid_video(slides, image_paths, audio_map, diagrams, work_dir)

        out_video = os.path.join(OUTPUT_DIR, "lecture.mp4")
        shutil.copy(video_path, out_video)

        out_meta = os.path.join(OUTPUT_DIR, "lecture_metadata.json")
        with open(out_meta, "w") as f:
            json.dump(structure, f, indent=2)

        if "full_description" in meta:
            out_desc = os.path.join(OUTPUT_DIR, "lecture_description.md")
            with open(out_desc, "w") as f:
                f.write(f"# {topic or 'Lecture'}\n\n{meta['full_description']}")

    logger.info("=" * 70)
    logger.info("DONE -> outputs/lecture.mp4")
    logger.info("=" * 70)
    return out_video


# ════════════════════════════════════════════════════════════════════════════
# SHARED VOICEOVER ENRICHMENT
# ════════════════════════════════════════════════════════════════════════════

def _build_source_block(
    transcript_text: str = None,
    notes_text: str = None,
    context: str = None,
    topic: str = "",
) -> tuple:
    ctx_parts: List[str] = []
    if topic and topic.strip():
        ctx_parts.append(f"COURSE TOPIC: {topic.strip()}")
    if context and context.strip():
        ctx_parts.append(f"COURSE CONTEXT:\n{context.strip()}")
    if transcript_text and transcript_text.strip():
        ctx_parts.append(f"CLASS TRANSCRIPT (primary source):\n{transcript_text[:20000]}")
    if notes_text and notes_text.strip():
        ctx_parts.append(f"NOTES / TEXTBOOK (depth reference):\n{notes_text[:8000]}")

    if ctx_parts:
        source_block = "\n\n---\n\n".join(ctx_parts)
        source_instruction = (
            "Use the source material below as the basis for your narration. "
            "Prioritise the transcript. Do NOT invent facts not in the source material."
        )
    else:
        source_block = "(No reference material provided.)"
        source_instruction = (
            "No reference material has been uploaded. "
            "Write rich, accurate academic narration based on your knowledge of the topic. "
            "Explain concepts clearly and thoroughly as a knowledgeable university lecturer would."
        )
    return source_block, source_instruction


def _enrich_voiceovers(
    slides: List[Dict],
    transcript_text: str = None,
    notes_text: str = None,
    context: str = None,
    topic: str = "",
) -> List[Dict]:
    BATCH = 8
    n = len(slides)
    source_block, source_instruction = _build_source_block(
        transcript_text, notes_text, context, topic
    )

    def _slide_context(slide: Dict) -> str:
        title   = slide.get("title", "Untitled Slide")
        bullets = slide.get("bullets", [])
        tcont   = slide.get("text_content", "").strip()
        lines   = [f"TITLE: {title}"]
        if bullets:
            lines.append("KEY POINTS:")
            for b in bullets:
                lines.append(f"  - {str(b).strip()}")
        if tcont and tcont != title:
            lines.append(f"SLIDE TEXT: {tcont[:600]}")
        return "\n".join(lines)

    def _call_batch(batch_slides: List[Dict]) -> List[Optional[str]]:
        slides_block = ""
        for local_i, sl in enumerate(batch_slides):
            slides_block += f"\n=== SLIDE {local_i + 1} ===\n{_slide_context(sl)}\n"

        prompt = f"""You are a university lecturer recording spoken audio narration for a lecture video.

{source_instruction}

SOURCE MATERIAL:
{source_block}

SLIDES TO NARRATE ({len(batch_slides)} slides):
{slides_block}

STRICT REQUIREMENTS:
1. Write narration for EVERY slide above, in EXACT ORDER (position 1, 2, 3...).
2. CONTENT / TECHNICAL slides: 120-180 words of fluent, engaging spoken narration.
   - EXPLAIN the concept in depth: what it is, why it matters, how it works.
   - INCORPORATE specific facts, terms, examples from the slide's key points.
   - Add ONE concrete real-world example or analogy.
   - Do NOT just read out bullet points. Do NOT produce placeholder text.
   - Do NOT start with "In this slide..." or "Today we will see...".
3. TITLE / INTRO / AGENDA / THANK-YOU slides: 15-30 words only.
4. Natural flowing spoken English. Conversational but authoritative.
5. TTS FORMATTING — ABSOLUTELY CRITICAL:
   - NO markdown whatsoever: no **bold**, no *italic*, no `backticks`, no ## headers.
   - NO bullet points or numbered lists in the voiceover text.
   - NO underscores for subscripts — write "b sub n" not "b_n", "x sub i" not "x_i".
   - NO carets for superscripts — write "x squared" not "x^2", "e to the x" not "e^x".
   - Spell out ALL math symbols in plain English: "the integral from zero to infinity",
     "alpha sub zero", "partial derivative of f with respect to x".
   - This text goes DIRECTLY to a text-to-speech voice engine. Plain English only.
6. Return STRICT JSON only — no markdown fences, no preamble, no trailing text.

JSON FORMAT (exactly {len(batch_slides)} objects, in order):
{{
  "slides": [
    {{"position": 1, "voiceover": "full narration for slide 1 here"}},
    {{"position": 2, "voiceover": "full narration for slide 2 here"}}
  ]
}}"""

        response = client.models.generate_content(
            model="gemini-2.5-flash", contents=prompt
        )
        raw  = response.text.strip()
        text = _repair_json(raw)
        try:
            data = json.loads(text)
        except json.JSONDecodeError:
            text = re.sub(r"(?<!\\)\\(?![\\\"/nrtbf])", r"\\\\", text)
            data = json.loads(text)

        items = data.get("slides", []) if isinstance(data, dict) else []

        pos_map: Dict[int, str] = {}
        for item in items:
            if isinstance(item, dict):
                vo = str(item.get("voiceover", "")).strip()
                p  = item.get("position")
                if p is not None and vo:
                    pos_map[int(p)] = vo

        result_vos: List[Optional[str]] = []
        for local_i in range(len(batch_slides)):
            vo = pos_map.get(local_i + 1)
            if not vo and local_i < len(items) and isinstance(items[local_i], dict):
                vo = str(items[local_i].get("voiceover", "")).strip() or None
            result_vos.append(vo)
        return result_vos

    def _single_fallback(slide: Dict) -> str:
        ctx    = _slide_context(slide)
        prompt = f"""You are a university lecturer. Record 120-160 words of spoken audio narration
for the following lecture slide.

{source_instruction}

SOURCE MATERIAL:
{source_block}

SLIDE:
{ctx}

Requirements:
- Explain the concept in depth: define it, explain why it matters, give a concrete example.
- Natural spoken English. Do NOT start with "In this slide..." or just restate the title.
- TTS FORMATTING — CRITICAL: No markdown (**bold**, *italic*, `code`, ## headers).
  No underscores (write "b sub n" not "b_n"). No carets (write "x squared" not "x^2").
  Plain spoken English only — this goes directly to a text-to-speech engine.
- Return ONLY the narration text. No JSON, no quotes, no preamble.
"""
        try:
            response = client.models.generate_content(
                model="gemini-2.5-flash", contents=prompt
            )
            vo = response.text.strip().strip("\"'`")
            return vo if len(vo.split()) >= 30 else ""
        except Exception as e:
            logger.warning(f"  Single fallback failed: {e}")
            return ""

    needs_vo: List[int] = [
        i for i, sl in enumerate(slides)
        if not str(sl.get("voiceover", "")).strip()
        or len(str(sl.get("voiceover", "")).split()) < 10
    ]
    if not needs_vo:
        logger.info("All slides already have voiceovers — skipping enrichment")
        return slides

    logger.info(f"Generating voiceovers for {len(needs_vo)}/{n} slides…")

    vo_results: Dict[int, Optional[str]] = {i: None for i in needs_vo}

    for batch_start in range(0, len(needs_vo), BATCH):
        batch_indices = needs_vo[batch_start: batch_start + BATCH]
        batch_slides  = [slides[i] for i in batch_indices]
        b_num = batch_start // BATCH + 1
        logger.info(f"  Voiceover batch {b_num}: slides {[i+1 for i in batch_indices]}…")
        try:
            vos = _call_batch(batch_slides)
            for local_i, vo in enumerate(vos):
                vo_results[batch_indices[local_i]] = vo
        except Exception as e:
            logger.warning(f"  Batch {b_num} failed: {e}")

    empty = [i for i in needs_vo if not vo_results.get(i)]
    if empty:
        logger.info(f"  {len(empty)} slide(s) need individual fallback: {[i+1 for i in empty]}")
    for i in empty:
        title = slides[i].get("title", f"Slide {i+1}")
        logger.info(f"    → Single call: \"{title[:55]}\"")
        vo = _single_fallback(slides[i])
        vo_results[i] = vo or None

    for i, vo in vo_results.items():
        if not vo:
            title       = slides[i].get("title", f"Slide {i+1}")
            bullets     = slides[i].get("bullets", [])
            bullet_text = " ".join(str(b) for b in bullets[:3])
            vo = (
                f"Let us now look at {title}. "
                + (f"{bullet_text}. " if bullet_text else "")
                + "This is an important concept in the course."
            )
            logger.warning(f"  Slide {i+1}: last-resort voiceover used")

        slides[i]["voiceover"] = vo
        wc = max(1, len(vo.split()))
        slides[i]["slide_duration_seconds"] = max(6, min(120, int(wc / 2.4)))

    n_good = sum(1 for i in needs_vo if len(slides[i].get("voiceover","").split()) > 30)
    logger.info(f"Voiceover enrichment done: {n_good}/{len(needs_vo)} with full content")
    return slides


# ════════════════════════════════════════════════════════════════════════════
# SLIDES-AS-IS PIPELINE
# ════════════════════════════════════════════════════════════════════════════

def extract_slide_titles_from_pptx(pptx_path: str) -> List[str]:
    return [s["title"] for s in extract_slide_content_from_pptx(pptx_path)]


def extract_slide_content_from_pptx(pptx_path: str) -> List[Dict]:
    prs = Presentation(pptx_path)
    result: List[Dict] = []

    for idx, slide in enumerate(prs.slides):
        title = ""
        body_texts: List[str] = []

        for shape in slide.shapes:
            if not getattr(shape, "has_text_frame", False):
                continue
            raw = shape.text.strip()
            if not raw:
                continue

            is_title = (shape == slide.shapes.title)
            lines = [l.strip() for l in raw.splitlines() if l.strip()]

            if is_title and lines:
                title = lines[0]
                if len(lines) > 1:
                    body_texts.extend(lines[1:])
            else:
                body_texts.extend(lines)

        if not title and body_texts:
            title = body_texts.pop(0)

        title = title or f"Slide {idx + 1}"
        text_content = " ".join(body_texts)
        result.append({"title": title, "text_content": text_content})

    return result


def generate_voiceovers_for_slide_titles(
    slide_data,
    transcript_text: str = None,
    notes_text: str = None,
    context: str = None,
) -> List[Dict]:
    if slide_data and isinstance(slide_data[0], str):
        slides = [{"title": t, "text_content": ""} for t in slide_data]
    else:
        slides = [dict(s) for s in slide_data]

    for s in slides:
        s["voiceover"] = ""

    enriched = _enrich_voiceovers(
        slides          = slides,
        transcript_text = transcript_text,
        notes_text      = notes_text,
        context         = context,
        topic           = "",
    )

    result: List[Dict] = []
    for idx, sl in enumerate(enriched):
        result.append({
            "slide_index":            idx + 1,
            "title":                  sl.get("title", f"Slide {idx + 1}"),
            "voiceover":              sl.get("voiceover", ""),
            "slide_duration_seconds": sl.get("slide_duration_seconds", 60),
        })
    return result


def _compose_slides_video(
    slides: List[Dict],
    image_paths: Dict[int, str],
    audio_map: Dict[int, Dict],
    work_dir: str,
) -> str:
    logger.info("Composing slides video…")
    TARGET_FPS = 30
    clip_files: List[str] = []

    for i, slide in enumerate(slides):
        try:
            image_path = image_paths.get(i)
            if not image_path or not os.path.exists(image_path):
                logger.warning(f"Slide {i + 1}: missing image, skipping")
                continue

            audio_info = audio_map.get(i, {})
            audio_path = audio_info.get("path", "")
            audio_ok   = bool(audio_path and os.path.exists(audio_path))
            duration   = float(
                audio_info.get("duration", slide.get("slide_duration_seconds", 8))
                if audio_ok
                else slide.get("slide_duration_seconds", 8)
            )

            clip_path = os.path.join(work_dir, f"clip_{i}.mp4")
            cmd = ["ffmpeg", "-y", "-loglevel", "error", "-loop", "1", "-i", image_path]
            if audio_ok:
                cmd += ["-i", audio_path]
            else:
                cmd += ["-f", "lavfi", "-i", "anullsrc=r=44100:cl=stereo"]
            cmd += [
                "-vf",
                "scale=1920:1080:force_original_aspect_ratio=decrease,"
                "pad=1920:1080:(ow-iw)/2:(oh-ih)/2:color=black",
                "-r", str(TARGET_FPS),
                "-c:v", "libx264", "-preset", "medium", "-crf", "20", "-pix_fmt", "yuv420p",
                "-c:a", "aac", "-b:a", "192k", "-ar", "44100", "-ac", "2",
                "-t", f"{duration:.2f}", "-shortest",
                clip_path,
            ]
            subprocess.run(cmd, check=True, capture_output=True, timeout=300)

            if os.path.exists(clip_path) and os.path.getsize(clip_path) > 5000:
                clip_files.append(clip_path)
                logger.info(f"Clip {i + 1}: {duration:.1f}s")
            else:
                logger.warning(f"Clip {i + 1}: output empty or missing")

        except Exception as e:
            logger.error(f"Slide {i + 1} clip error: {e}")

    if not clip_files:
        raise ValueError("No clips were created!")

    concat_file = os.path.join(work_dir, "concat.txt")
    with open(concat_file, "w") as f:
        for clip in clip_files:
            f.write(f"file '{clip}'\n")

    output = os.path.join(work_dir, "final.mp4")
    subprocess.run(
        ["ffmpeg", "-y", "-f", "concat", "-safe", "0", "-i", concat_file, "-c", "copy", output],
        check=True, capture_output=True, timeout=600,
    )
    if not os.path.exists(output):
        raise FileNotFoundError("Final slides video was not created")

    logger.info(f"Slides video: {os.path.getsize(output) / (1024 * 1024):.1f} MB")
    return output


def create_slides_lecture(
    pptx_path: str,
    transcript_text: str = None,
    notes_text: str = None,
    context: str = None,
    slides_override: List[Dict] = None,
) -> str:
    logger.info("=" * 70)
    logger.info("Slides-as-is pipeline starting…")
    logger.info("=" * 70)

    with tempfile.TemporaryDirectory() as work_dir:
        image_paths = pptx_to_images(pptx_path, work_dir)
        if not image_paths:
            raise RuntimeError(
                "PPTX → image conversion produced no output. "
                "Make sure LibreOffice (soffice) is installed and in PATH."
            )
        logger.info(f"Slide images ready: {len(image_paths)}")

        if slides_override is not None:
            slides = slides_override
            logger.info(f"Using {len(slides)} user-reviewed voiceover scripts")
        else:
            slide_contents = extract_slide_content_from_pptx(pptx_path)
            slides = generate_voiceovers_for_slide_titles(
                slide_contents, transcript_text, notes_text, context
            )

        audio_map  = generate_voiceovers(slides, work_dir)
        video_path = _compose_slides_video(slides, image_paths, audio_map, work_dir)

        out_video = os.path.join(OUTPUT_DIR, "lecture.mp4")
        shutil.copy(video_path, out_video)

        total_dur = sum(s.get("slide_duration_seconds", 8) for s in slides)
        meta = {
            "total_slides": len(slides),
            "estimated_duration_minutes": max(1, total_dur // 60),
            "target_audience": "intermediate",
            "full_description": (
                "Narrated lecture video created from an uploaded PPTX. "
                "Slides are unchanged; only AI-generated voiceover narration was added."
            ),
        }
        with open(os.path.join(OUTPUT_DIR, "lecture_metadata.json"), "w") as f:
            json.dump({"metadata": meta, "slides": slides}, f, indent=2)

    logger.info("=" * 70)
    logger.info("Slides pipeline DONE → outputs/lecture.mp4")
    logger.info("=" * 70)
    return out_video


# ════════════════════════════════════════════════════════════════════════════
# ══ TOPIC VIDEO PIPELINE (Pure Manim) ══════════════════════════════════════
# ════════════════════════════════════════════════════════════════════════════

_TOPIC_STRUCTURE_PROMPT = """You are an expert educational video scriptwriter AND Manim animation designer.

Create a structured scene plan for a short educational video on the given topic.
Each scene maps to one Manim animation using our template library.

TOPIC: {topic}
CONTEXT: {context}
TARGET DURATION: {target_minutes} minutes (~{target_seconds} seconds total)

AVAILABLE SCENE TYPES (choose the best for each scene):
{scene_types_doc}

OUTPUT FORMAT — VALID JSON ONLY, NOTHING ELSE:
{{
  "video_title": "concise video title",
  "description": "2-3 sentence summary",
  "total_scenes": integer,
  "estimated_minutes": number,
  "scenes": [
    {{
      "scene_index": 0,
      "title": "scene title (max 50 chars, ASCII)",
      "scene_type": "one of the available scene types",
      "duration_seconds": integer (8-90),
      "voiceover": "spoken narration for this scene (120-180 words for content scenes, 15-30 for title/summary)",
      "params": {{
        "title": "same as scene title",
        ... scene-type-specific fields ...
      }}
    }}
  ]
}}

SCENE PLANNING RULES:
1. ALWAYS start with a title_card scene (10-12 seconds).
2. ALWAYS end with a summary_card scene (15-20 seconds).
3. Build logical flow: intro → concepts → examples → diagrams → summary.
4. Use math_steps or worked_example for quantitative content.
5. Use definition_box for introducing key terms.
6. Use bullet_list for conceptual explanations with 3-5 key points.
7. Use calculus_plot/bar_chart/line_graph for data/function visualization.
8. Use concept_map for showing relationships between ideas.
9. Vary scene types — don't use the same type more than 3 times in a row.
10. Total duration should be close to {target_seconds} seconds.
11. voiceover MUST be substantive — explain WHY and HOW, not just WHAT.
12. All strings: ASCII only.
13. Output ONLY valid JSON. No markdown, no explanation.

For each scene_type, here are the required params fields:
title_card:       title (str), subtitle (str), color (hex)
bullet_list:      title (str), bullets (list of 3-5 strings), color (hex)
definition_box:   title (str), term (str), definition (str), example (str), color (hex)
math_steps:       title (str), steps: [{{eq: latex, note: str}}]
worked_example:   title (str), problem (latex str), solution_steps: [{{eq: latex, note: str}}]
formula:          title (str), latex (str), explanation (str)
calculus_plot:    title (str), func_latex (str), func_expr (python), x_range ([min,max]), y_range ([min,max])
bar_chart:        title (str), labels (list), values (list of floats), y_label (str)
line_graph:       title (str), curves: [{{label, values: [floats]}}], x_label, y_label
pie_chart:        title (str), slices: [{{label, value}}]
block_flow:       title (str), boxes (list of 2-6 str)
neural_network:   title (str), layer_sizes (list of ints), layer_labels (list of str)
concept_map:      title (str), center (str), branches: [{{label: str}}]
timeline:         title (str), events: [{{label, year}}]
step_by_step:     title (str), steps (list of 3-6 strings)
comparison_table: title (str), col_a (str), col_b (str), rows: [{{a: str, b: str}}]
venn_diagram:     title (str), circle_a (str), circle_b (str), intersection (str), items_a (list), items_b (list)
gradient_descent: title (str), steps (int 4-8), learning_rate (float 0.1-0.5)
physics_forces:   title (str), body (str), forces: [{{label, dx, dy, color}}]
tree_diagram:     title (str), root (str), children (list of str)
matrix_display:   title (str), label (str), matrix_latex (latex str), note (str)
probability_dist: title (str), mean (float), std (float)
summary_card:     title (str), points (list of 3-5 key takeaways), color (hex)
"""


def generate_topic_structure(
    topic: str,
    context: str = "",
    target_minutes: int = 7,
) -> Dict:
    """Generate a complete scene plan for a pure Manim topic video."""
    logger.info(f"Generating topic video structure: '{topic}' (~{target_minutes} min)")

    target_seconds = target_minutes * 60
    scene_types_list = sorted(AVAILABLE_SCENE_TYPES)
    scene_types_doc  = "\n".join(f"  - {t}" for t in scene_types_list)

    prompt = _TOPIC_STRUCTURE_PROMPT.format(
        topic          = topic,
        context        = context or "No additional context provided.",
        target_minutes = target_minutes,
        target_seconds = target_seconds,
        scene_types_doc = scene_types_doc,
    )

    structure = None
    for attempt in range(3):
        try:
            response = client.models.generate_content(model="gemini-2.5-flash", contents=prompt)
            raw  = response.text.strip()
            text = _repair_json(raw)
            structure = json.loads(text)
            logger.info(
                f"Topic structure OK (attempt {attempt+1}): "
                f"{structure.get('total_scenes','?')} scenes, "
                f"~{structure.get('estimated_minutes','?')} min"
            )
            break
        except (json.JSONDecodeError, Exception) as e:
            logger.warning(f"Topic structure attempt {attempt+1}/3 failed: {e}")
            if attempt == 2:
                raise RuntimeError(f"Failed to generate topic structure after 3 attempts: {e}") from e
            time.sleep(2)

    return structure


def render_topic_scenes(
    scenes: List[Dict],
    work_dir: str,
    topic: str = "",
) -> Dict[int, str]:
    """Render all scenes using the template library. Returns {idx: mp4_path}."""
    if not MANIM_BIN or not os.path.exists(MANIM_BIN):
        logger.warning(f"Manim not found at '{MANIM_BIN}'")
        return {}

    env = os.environ.copy()
    env["PATH"] = "/Library/TeX/texbin:" + env.get("PATH", "")

    try:
        r = subprocess.run([MANIM_BIN, "--version"],
                           capture_output=True, text=True, timeout=120, env=env)
        if r.returncode != 0:
            logger.warning(f"Manim check failed: {r.stderr[:200]}")
            return {}
        logger.info(f"Manim ready: {r.stdout.strip()}")
    except Exception as e:
        logger.warning(f"Manim unavailable: {e}")
        return {}

    logger.info(f"Rendering {len(scenes)} topic scenes…")
    outputs: Dict[int, str] = {}

    def render_one(scene: Dict) -> Tuple[int, Optional[str]]:
        i        = scene["scene_index"]
        stype    = scene.get("scene_type", "bullet_list")
        duration = int(scene.get("duration_seconds", 40))
        params   = dict(scene.get("params") or {})
        params["type"]  = stype
        params["title"] = params.get("title") or scene.get("title", f"Scene {i+1}")

        scene_dir = os.path.join(work_dir, f"topic_{i}")
        os.makedirs(scene_dir, exist_ok=True)

        code = get_scene_code(i, params, duration)
        if not code:
            logger.warning(f"  Scene {i+1}: no code generated for type '{stype}'")
            return i, None

        ok, err = _validate_code(code, i)
        if not ok:
            logger.error(f"  Scene {i+1}: invalid code: {err}")
            # Try fallback to bullet_list
            fallback_params = {
                "type":    "bullet_list",
                "title":   params.get("title", f"Scene {i+1}"),
                "bullets": [
                    params.get("title", f"Scene {i+1}"),
                    "See course materials for details.",
                ],
            }
            code = get_scene_code(i, fallback_params, duration)
            if not code:
                return i, None
            ok, err = _validate_code(code, i)
            if not ok:
                return i, None

        video, render_err = _try_render(code, i, scene_dir, env)
        if video:
            logger.info(f"  Scene {i+1} ({stype}): rendered OK")
        else:
            logger.warning(f"  Scene {i+1} ({stype}): failed — {(render_err or '')[:120]}")
        return i, video

    with ThreadPoolExecutor(max_workers=2) as ex:
        futures = {ex.submit(render_one, sc): sc["scene_index"] for sc in scenes}
        for future in as_completed(futures):
            i, path = future.result()
            if path:
                outputs[i] = path

    logger.info(f"Topic scenes rendered: {len(outputs)}/{len(scenes)}")
    return outputs


def _compose_topic_video(
    scenes: List[Dict],
    rendered_scenes: Dict[int, str],
    audio_map: Dict[int, Dict],
    work_dir: str,
) -> str:
    """Compose topic video: Manim scenes + TTS audio → final MP4."""
    logger.info("Composing topic video…")
    TARGET_FPS = 30
    clip_files: List[str] = []

    for scene in scenes:
        i = scene["scene_index"]
        try:
            video_src  = rendered_scenes.get(i)
            audio_info = audio_map.get(i, {})
            audio_path = audio_info.get("path", "")
            audio_ok   = bool(audio_path and os.path.exists(audio_path))
            duration   = float(
                audio_info.get("duration", scene.get("duration_seconds", 30))
                if audio_ok
                else scene.get("duration_seconds", 30)
            )

            clip_path = os.path.join(work_dir, f"topicclip_{i}.mp4")

            if video_src and os.path.exists(video_src):
                # Manim video + TTS audio
                cmd = ["ffmpeg", "-y", "-loglevel", "error",
                       "-stream_loop", "-1", "-i", video_src]
                if audio_ok:
                    cmd += ["-i", audio_path]
                else:
                    cmd += ["-f", "lavfi", "-i", "anullsrc=r=44100:cl=stereo"]
                cmd += [
                    "-vf", "scale=1920:1080:force_original_aspect_ratio=decrease,"
                           "pad=1920:1080:(ow-iw)/2:(oh-ih)/2:color=#0F172A",
                    "-r", str(TARGET_FPS),
                    "-c:v", "libx264", "-preset", "medium", "-crf", "20", "-pix_fmt", "yuv420p",
                    "-c:a", "aac", "-b:a", "192k", "-ar", "44100", "-ac", "2",
                    "-t", f"{duration:.2f}", "-shortest",
                    clip_path,
                ]
            else:
                # Fallback: plain dark frame with title text overlay + audio
                title_safe = re.sub(r"[^\w\s]", "", scene.get("title", f"Scene {i+1}"))[:50]
                cmd = ["ffmpeg", "-y", "-loglevel", "error",
                       "-f", "lavfi", "-i", f"color=c=0x0F172A:size=1920x1080:rate={TARGET_FPS}"]
                if audio_ok:
                    cmd += ["-i", audio_path]
                else:
                    cmd += ["-f", "lavfi", "-i", "anullsrc=r=44100:cl=stereo"]
                cmd += [
                    "-vf", f"drawtext=text='{title_safe}':fontsize=48:fontcolor=0x00D9FF:"
                           f"x=(w-text_w)/2:y=(h-text_h)/2",
                    "-c:v", "libx264", "-preset", "medium", "-crf", "20", "-pix_fmt", "yuv420p",
                    "-c:a", "aac", "-b:a", "192k", "-ar", "44100", "-ac", "2",
                    "-t", f"{duration:.2f}", "-shortest",
                    clip_path,
                ]

            subprocess.run(cmd, check=True, capture_output=True, timeout=300)

            if os.path.exists(clip_path) and os.path.getsize(clip_path) > 5000:
                clip_files.append(clip_path)
                logger.info(f"Topic clip {i+1}: {duration:.1f}s")
            else:
                logger.warning(f"Topic clip {i+1}: empty/missing")

        except Exception as e:
            logger.error(f"Topic scene {i+1} composition error: {e}")

    if not clip_files:
        raise ValueError("No topic clips created!")

    concat_file = os.path.join(work_dir, "topic_concat.txt")
    with open(concat_file, "w") as f:
        for clip in clip_files:
            f.write(f"file '{clip}'\n")

    output = os.path.join(work_dir, "topic_final.mp4")
    subprocess.run(
        ["ffmpeg", "-y", "-f", "concat", "-safe", "0", "-i", concat_file, "-c", "copy", output],
        check=True, capture_output=True, timeout=600,
    )
    if not os.path.exists(output):
        raise FileNotFoundError("Final topic video not created")

    size_mb = os.path.getsize(output) / (1024 * 1024)
    logger.info(f"Topic video: {size_mb:.1f} MB")
    return output


def create_topic_video(
    topic: str,
    context: str = "",
    target_minutes: int = 7,
    scenes_override: List[Dict] = None,
) -> str:
    """
    Pure Manim Topic Video Pipeline.
    
    1. Gemini generates a scene plan (JSON).
    2. Each scene is rendered via Manim template library.
    3. TTS voiceover is generated with Edge TTS.
    4. Scenes + audio are composed into a final MP4.
    
    Returns path to the final MP4 in outputs/.
    """
    logger.info("=" * 70)
    logger.info(f"Topic Video Pipeline: '{topic}'")
    logger.info("=" * 70)

    with tempfile.TemporaryDirectory() as work_dir:
        # 1. Get scene plan
        if scenes_override is not None:
            scenes = scenes_override
            logger.info(f"Using {len(scenes)} user-reviewed scenes")
            video_title = topic
        else:
            structure   = generate_topic_structure(topic, context, target_minutes)
            scenes      = structure.get("scenes", [])
            video_title = structure.get("video_title", topic)

        if not scenes:
            raise ValueError("No scenes generated — cannot create video.")

        logger.info(f"Scenes to render: {len(scenes)}")

        # 2. Render Manim scenes + generate audio in parallel
        def _render_all():
            return render_topic_scenes(scenes, work_dir, topic=topic)

        def _audio_all():
            return generate_voiceovers(scenes, work_dir)

        with ThreadPoolExecutor(max_workers=2) as ex:
            f_render = ex.submit(_render_all)
            f_audio  = ex.submit(_audio_all)
            rendered_scenes = f_render.result()
            audio_map       = f_audio.result()

        logger.info(f"Rendered: {len(rendered_scenes)}/{len(scenes)} scenes")

        # 3. Compose
        video_path = _compose_topic_video(scenes, rendered_scenes, audio_map, work_dir)

        # 4. Save outputs
        out_video = os.path.join(OUTPUT_DIR, "lecture.mp4")
        shutil.copy(video_path, out_video)

        total_dur = sum(s.get("duration_seconds", 30) for s in scenes)
        meta = {
            "total_slides":               len(scenes),
            "estimated_duration_minutes": max(1, total_dur // 60),
            "target_audience":            "intermediate",
            "full_description":           (
                f"Pure Manim educational video on '{topic}'. "
                f"{len(scenes)} animated scenes rendered with the Manim Community library. "
                f"Voiceover generated with Edge TTS."
            ),
        }
        with open(os.path.join(OUTPUT_DIR, "lecture_metadata.json"), "w") as f:
            json.dump({"metadata": meta, "scenes": scenes, "video_title": video_title}, f, indent=2)

        out_desc = os.path.join(OUTPUT_DIR, "lecture_description.md")
        with open(out_desc, "w") as f:
            f.write(f"# {video_title}\n\nPure Manim animated video on: {topic}\n")

    logger.info("=" * 70)
    logger.info("Topic Video DONE → outputs/lecture.mp4")
    logger.info("=" * 70)
    return out_video


def regenerate_topic_scene(
    scene_idx: int,
    current_scene: Dict,
    topic: str,
    context: str = "",
    surrounding_scenes: List[Dict] = None,
    custom_prompt: str = "",
) -> Dict:
    """Regenerate a single scene in the topic video pipeline."""
    logger.info(f"Regenerating topic scene {scene_idx + 1}")
    surrounding_scenes = surrounding_scenes or []

    surrounding_info = "\n".join(
        f"  Scene {s.get('scene_index','?')}: \"{s.get('title','?')}\" [{s.get('scene_type','?')}]"
        for s in surrounding_scenes
    ) or "  (none)"

    custom_section = (
        f"\nSPECIAL INSTRUCTIONS (highest priority):\n{custom_prompt.strip()}\n"
        if custom_prompt and custom_prompt.strip() else ""
    )

    scene_types_list = sorted(AVAILABLE_SCENE_TYPES)
    scene_types_doc  = "\n".join(f"  - {t}" for t in scene_types_list)

    prompt = f"""You are regenerating ONE scene in a pure Manim educational video.

TOPIC: {topic}
CONTEXT: {context or 'None'}

SURROUNDING SCENES:
{surrounding_info}

CURRENT SCENE:
{json.dumps(current_scene, indent=2)}

{custom_section}

AVAILABLE SCENE TYPES:
{scene_types_doc}

Generate a new version of this scene. Return ONLY a valid JSON object with these fields:
scene_index, title, scene_type, duration_seconds, voiceover, params

The params must match the chosen scene_type exactly.
No markdown, no explanation.
"""
    response = client.models.generate_content(model="gemini-2.5-flash", contents=prompt)
    text = _repair_json(response.text.strip())
    try:
        result = json.loads(text)
    except json.JSONDecodeError:
        text = re.sub(r"(?<!\\)\\(?![\\\"nrtbf/u])", r"\\\\", text)
        result = json.loads(text)

    result["scene_index"] = current_scene.get("scene_index", scene_idx)
    logger.info(f"Topic scene {scene_idx + 1} regenerated OK")
    return result


# ════════════════════════════════════════════════════════════════════════════
# CLI ENTRY POINT
# ════════════════════════════════════════════════════════════════════════════

if __name__ == "__main__":
    mode = input("Mode — (1) Full AI  (2) Topic Video: ").strip()
    if mode == "2":
        topic   = input("Topic: ").strip()
        minutes = int(input("Target minutes (5-10): ").strip() or "7")
        create_topic_video(topic=topic, target_minutes=minutes)
    else:
        topic = input("Topic (or Enter to infer): ").strip()
        ref   = load_reference_documents_from_project()
        create_lecture(topic=topic, notes_text=ref if ref.strip() else None)
