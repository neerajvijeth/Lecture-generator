import streamlit as st
import tempfile
import os
import time
import logging
import json
import traceback
import copy
from pathlib import Path

from pptx import Presentation
from docx import Document

try:
    from app import (
        create_lecture,
        load_pdf_text,
        regenerate_single_slide,
        generate_lecture_structure,
        fetch_preview_image,
        invalidate_preview_cache,
        # Slides-as-is pipeline
        extract_slide_titles_from_pptx,
        extract_slide_content_from_pptx,
        generate_voiceovers_for_slide_titles,
        create_slides_lecture,
        # Topic Video pipeline
        generate_topic_structure,
        create_topic_video,
        regenerate_topic_scene,
        render_topic_scenes,
        generate_voiceovers,
        OUTPUT_DIR,
    )
    from diagrams import AVAILABLE_DIAGRAM_TYPES, AVAILABLE_SCENE_TYPES
except Exception as e:
    st.error(f"Failed to import modules: {e}")
    st.stop()

# ════════════════════════════════════════════════════════════════════════════
# CONSTANTS
# ════════════════════════════════════════════════════════════════════════════

LOGO_PATH = "/Users/neerajvijeth/Desktop/gemini project/images.png"

# ════════════════════════════════════════════════════════════════════════════
# PAGE CONFIG
# ════════════════════════════════════════════════════════════════════════════
st.set_page_config(
    page_title="AI Lecture Generator",
    page_icon="🎓",
    layout="wide",
    initial_sidebar_state="expanded",
)

# ════════════════════════════════════════════════════════════════════════════
# CSS
# ════════════════════════════════════════════════════════════════════════════
st.markdown("""
<style>
  @import url('https://fonts.googleapis.com/css2?family=Inter:wght@400;500;600;700&display=swap');
  html, body, [class*="css"] { font-family: 'Inter', 'Segoe UI', sans-serif; }

  [data-testid="stSidebar"] { background: #080e1a !important; border-right: 1px solid #1e293b; }
  [data-testid="stSidebar"] * { color: #e2e8f0 !important; }
  [data-testid="stSidebar"] .stButton > button {
    background: linear-gradient(135deg, #0284c7, #00d9ff) !important;
    color: #ffffff !important; font-weight: 700; border: none;
    border-radius: 10px; letter-spacing: 0.02em;
  }

  .main .block-container { padding-top: 1.2rem; padding-bottom: 2rem; }

  .stButton > button[kind="primary"] {
    background: linear-gradient(135deg, #0284c7, #00d9ff) !important;
    color: #ffffff !important; font-weight: 700; border-radius: 10px;
    border: none; letter-spacing: 0.02em; transition: all 0.15s;
  }
  .stButton > button[kind="primary"]:hover { transform: translateY(-1px); box-shadow: 0 4px 15px rgba(0,217,255,0.3); }

  .stProgress > div > div { background: linear-gradient(90deg, #0284c7, #00d9ff); border-radius: 4px; }

  .stDownloadButton > button {
    background: #1e293b !important; color: #00d9ff !important;
    border: 1px solid #334155; border-radius: 10px; font-weight: 600;
  }

  .metric-card {
    background: linear-gradient(145deg, #1e293b, #162032);
    border-radius: 14px; padding: 20px 14px; text-align: center;
    border: 1px solid #2a3a52; min-height: 100px;
    display: flex; flex-direction: column; justify-content: center;
    transition: transform 0.15s;
  }
  .metric-card:hover { transform: translateY(-2px); border-color: #3d5470; }
  .metric-value { font-size: 1.5rem; font-weight: 700; color: #00d9ff; line-height: 1.2; }
  .metric-label { font-size: 0.72rem; color: #64748b; margin-top: 6px; letter-spacing: 0.04em; text-transform: uppercase; }

  .upload-header {
    background: linear-gradient(90deg, #162032, #1e293b);
    border-left: 3px solid #00d9ff;
    padding: 9px 14px; border-radius: 0 8px 8px 0;
    margin-bottom: 5px; font-size: 0.85rem; font-weight: 600; color: #e2e8f0;
  }
  .upload-hint { font-size: 0.72rem; color: #475569; margin-bottom: 10px; padding-left: 4px; }

  .scard {
    background: linear-gradient(145deg, #1a2540, #141e30);
    padding: 14px 18px 12px 18px; border-radius: 12px; margin-bottom: 6px;
    border: 1px solid #243048; transition: border-color 0.2s, box-shadow 0.2s;
  }
  .scard:hover { border-color: #3d5470; box-shadow: 0 2px 12px rgba(0,0,0,0.3); }
  .s-num   { color: #3d5470; font-size: 11px; font-weight: 700; letter-spacing: 0.08em; }
  .s-title { color: #e2e8f0; font-size: 15px; font-weight: 600; }
  .s-badge {
    display: inline-block; font-size: 10px; font-weight: 700;
    padding: 2px 9px; border-radius: 20px; margin-left: 7px; letter-spacing: 0.03em;
  }
  .b-text    { background: #0c2a4a; color: #60a5fa; border: 1px solid #1e4470; }
  .b-diagram { background: #1e1040; color: #a78bfa; border: 1px solid #3d2080; }
  .b-scene   { background: #0a2010; color: #4ade80; border: 1px solid #145214; }
  .b-orig    { background: #1a2030; color: #475569; border: 1px solid #2a3448; }
  .b-edited  { background: #082040; color: #38bdf8; border: 1px solid #1060a0; }
  .b-regen   { background: #1a0840; color: #c4b5fd; border: 1px solid #4020a0; }
  .b-dur     { color: #3d5470; font-size: 10px; margin-left: 6px; }

  .bullet-prev {
    font-size: 12px; color: #4a6080; padding: 3px 0 3px 14px;
    border-left: 2px solid #1e3048; margin: 3px 0;
    overflow: hidden; text-overflow: ellipsis; white-space: nowrap;
  }

  .vo-card {
    background: linear-gradient(145deg, #111e30, #0d1828);
    border: 1px solid #1e3048; border-radius: 10px;
    padding: 14px 18px; margin-bottom: 8px; transition: border-color 0.2s;
  }
  .vo-card:hover { border-color: #2a4060; }

  .img-placeholder {
    background: linear-gradient(145deg, #0d1525, #111e30);
    border: 1px dashed #2a3a52; border-radius: 10px;
    padding: 32px 12px; text-align: center;
    color: #3d5470; font-size: 11px; line-height: 1.7;
  }
  .img-badge {
    display: inline-flex; align-items: center; gap: 4px;
    font-size: 10px; font-weight: 700; padding: 3px 9px;
    border-radius: 20px; margin-bottom: 6px; letter-spacing: 0.03em;
  }
  .ib-custom  { background: #052010; color: #4ade80; border: 1px solid #145214; }
  .ib-fetched { background: #041830; color: #60a5fa; border: 1px solid #103060; }

  .log-box {
    background: #060d1a; border: 1px solid #1e293b; border-radius: 10px;
    padding: 12px 16px; font-family: 'SF Mono','Fira Code','Consolas',monospace;
    font-size: 11px; max-height: 300px; overflow-y: auto; color: #4a6080; line-height: 1.7;
  }
  .log-box::-webkit-scrollbar { width: 4px; }
  .log-box::-webkit-scrollbar-thumb { background: #1e3048; border-radius: 4px; }

  .panel {
    background: linear-gradient(145deg, #0d1525, #111e30);
    border: 1px solid #1e3048; border-radius: 10px;
    padding: 16px 18px; margin: 8px 0 12px 0;
  }
  .panel-title {
    font-size: 12px; font-weight: 700; color: #64748b;
    margin-bottom: 12px; letter-spacing: 0.08em; text-transform: uppercase;
  }

  .phase-banner {
    background: linear-gradient(135deg, #0d1830, #111e30);
    border: 1px solid #1e3048; border-radius: 12px;
    padding: 14px 20px; margin-bottom: 20px;
    font-size: 13px; color: #64748b; line-height: 1.6;
  }

  .scene-card {
    background: linear-gradient(145deg, #0d1f2a, #0a1820);
    border: 1px solid #1e3a2a; border-radius: 12px;
    padding: 14px 18px; margin-bottom: 8px;
    transition: border-color 0.2s, box-shadow 0.2s;
  }
  .scene-card:hover { border-color: #2a5a40; box-shadow: 0 2px 12px rgba(0,0,0,0.3); }
  .scene-type-badge {
    display: inline-block; font-size: 10px; font-weight: 700;
    padding: 2px 9px; border-radius: 20px; letter-spacing: 0.03em;
    background: #0a2010; color: #4ade80; border: 1px solid #145214;
    margin-left: 8px;
  }

  .section-sep {
    height: 1px; background: linear-gradient(90deg, transparent, #1e3048, transparent);
    margin: 20px 0;
  }

  hr { border-color: #1e293b !important; }

  .pipeline-pill {
    display: inline-block; padding: 4px 12px; border-radius: 20px;
    font-size: 11px; font-weight: 700; letter-spacing: 0.04em;
  }
  .pp-full   { background: #0c2a4a; color: #60a5fa; border: 1px solid #1e4470; }
  .pp-slides { background: #1e1040; color: #a78bfa; border: 1px solid #3d2080; }
  .pp-topic  { background: #0a2010; color: #4ade80; border: 1px solid #145214; }
</style>
""", unsafe_allow_html=True)


# ════════════════════════════════════════════════════════════════════════════
# SESSION STATE
# ════════════════════════════════════════════════════════════════════════════

def _init():
    defaults = {
        "phase":                "upload",
        "pipeline":             "full",
        "review_phase":         "review",
        # Full AI pipeline
        "structure":            None,
        "topic":                "",
        "transcript_text":      None,
        "slides_text":          None,
        "notes_text":           None,
        "skip_diagrams":        False,
        "slide_statuses":       {},
        "slide_edits":          {},
        "card_modes":           {},
        "regen_queue":          {},
        "slide_preview_images": {},
        "slide_custom_images":  {},
        "images_prefetched":    False,
        "pending_image_fetches": set(),
        # Slides-as-is pipeline
        "slides_pptx_bytes":    None,
        "slides_pptx_filename": "",
        "slides_pptx_path":     None,
        "slides_context":       "",
        "slides_structure":     None,
        "slides_edits":         {},
        # Topic Video pipeline
        "topic_topic":          "",
        "topic_context":        "",
        "topic_minutes":        7,
        "topic_structure":      None,
        "topic_scenes":         None,
        "topic_scene_edits":    {},
        "topic_scene_statuses": {},
        "topic_card_modes":     {},
        "topic_regen_queue":    {},
        # Shared output
        "output_path":          None,
        "render_elapsed":       None,
    }
    for k, v in defaults.items():
        if k not in st.session_state:
            st.session_state[k] = v

_init()


# ════════════════════════════════════════════════════════════════════════════
# UTILITY
# ════════════════════════════════════════════════════════════════════════════

def _read_image_bytes(path: str) -> bytes:
    with open(path, "rb") as f:
        return f.read()

def _fetch_and_store_preview(slide_idx: int, slide: dict, topic: str,
                              force_refresh: bool = False) -> bool:
    path = fetch_preview_image(slide_idx, slide, topic=topic, force_refresh=force_refresh)
    if path and os.path.exists(path) and os.path.getsize(path) > 100:
        st.session_state.slide_preview_images[slide_idx] = _read_image_bytes(path)
        return True
    st.session_state.slide_preview_images.pop(slide_idx, None)
    return False

def _show_logo_topright():
    if os.path.exists(LOGO_PATH):
        _, logo_col = st.columns([10, 1])
        with logo_col:
            st.image(LOGO_PATH, width=68)

def _show_logo_sidebar():
    if os.path.exists(LOGO_PATH):
        st.image(LOGO_PATH, width=86)
        st.markdown("<div style='margin-bottom:2px'></div>", unsafe_allow_html=True)

def _extract(files, label: str) -> str:
    texts = []
    for file in files:
        suffix   = Path(file.name).suffix.lower()
        tmp_path = None
        try:
            if suffix == ".pdf":
                with tempfile.NamedTemporaryFile(delete=False, suffix=".pdf") as tmp:
                    tmp.write(file.read()); tmp_path = tmp.name
                texts.append(load_pdf_text(tmp_path))
            elif suffix == ".docx":
                with tempfile.NamedTemporaryFile(delete=False, suffix=".docx") as tmp:
                    tmp.write(file.read()); tmp_path = tmp.name
                doc = Document(tmp_path)
                texts.extend(p.text for p in doc.paragraphs if p.text.strip())
            elif suffix == ".pptx":
                with tempfile.NamedTemporaryFile(delete=False, suffix=".pptx") as tmp:
                    tmp.write(file.read()); tmp_path = tmp.name
                prs = Presentation(tmp_path)
                for sl in prs.slides:
                    for shape in sl.shapes:
                        if shape.has_text_frame:
                            texts.append(shape.text)
            elif suffix in (".txt", ".md"):
                texts.append(file.read().decode("utf-8", errors="ignore"))
        except Exception as e:
            st.warning(f"Could not read {file.name} ({label}): {e}")
        finally:
            if tmp_path and os.path.exists(tmp_path):
                try: os.unlink(tmp_path)
                except Exception: pass
    return "\n\n".join(t for t in texts if t.strip())


# ════════════════════════════════════════════════════════════════════════════
# LOG HANDLER
# ════════════════════════════════════════════════════════════════════════════

class _LogHandler(logging.Handler):
    def __init__(self, buf):
        super().__init__()
        self.buf = buf
    def emit(self, record):
        self.buf.append(self.format(record))


# ════════════════════════════════════════════════════════════════════════════
# FULL PIPELINE HELPERS
# ════════════════════════════════════════════════════════════════════════════

def _get_slide(i: int) -> dict:
    return st.session_state.slide_edits.get(i, st.session_state.structure["slides"][i])

def _get_final_slides() -> list:
    return [_get_slide(i) for i in range(len(st.session_state.structure["slides"]))]

def _status(i: int) -> str:
    return st.session_state.slide_statuses.get(i, "original")

def _n_modified() -> int:
    return sum(1 for v in st.session_state.slide_statuses.values() if v != "original")

def _open_edit(i: int):
    sl = _get_slide(i)
    st.session_state.card_modes[i] = "edit"
    st.session_state[f"et_{i}"]  = sl.get("title", "")
    st.session_state[f"eb_{i}"]  = "\n".join(sl.get("bullets", []))
    st.session_state[f"ev_{i}"]  = sl.get("voiceover", "")
    st.session_state[f"ety_{i}"] = sl.get("type", "text")
    diag = sl.get("diagram", "none") or "none"
    opts = ["none"] + sorted(AVAILABLE_DIAGRAM_TYPES)
    st.session_state[f"ed_{i}"]  = diag if diag in opts else "none"
    img = sl.get("image") or {}
    st.session_state[f"ein_{i}"] = img.get("needed", False)
    st.session_state[f"eiq_{i}"] = "\n".join(img.get("search_queries", []))

def _save_edit(i: int):
    original = st.session_state.structure["slides"][i]
    bullets  = [b.strip() for b in st.session_state.get(f"eb_{i}", "").split("\n") if b.strip()]
    queries  = [q.strip() for q in st.session_state.get(f"eiq_{i}", "").split("\n") if q.strip()]
    edited   = copy.deepcopy(original)
    edited["title"]     = st.session_state.get(f"et_{i}", original.get("title", ""))[:55]
    edited["bullets"]   = bullets
    edited["voiceover"] = st.session_state.get(f"ev_{i}", original.get("voiceover", ""))
    new_type = st.session_state.get(f"ety_{i}", original.get("type", "text"))
    edited["type"]      = new_type
    diag = st.session_state.get(f"ed_{i}", "none")
    edited["diagram"]        = diag if diag != "none" else "none"
    edited["diagram_needed"] = diag != "none"
    edited["image"] = {
        "needed":         st.session_state.get(f"ein_{i}", False),
        "search_queries": queries,
        "placement":      "right",
    }
    st.session_state.slide_edits[i]    = edited
    st.session_state.slide_statuses[i] = "edited"
    st.session_state.card_modes[i]     = None
    _clear_preview_bytes(i)
    if new_type == "text" and edited["image"]["needed"] and queries:
        st.session_state.pending_image_fetches.add(i)
    elif new_type == "diagram":
        st.session_state.slide_preview_images.pop(i, None)
        st.session_state.slide_custom_images.pop(i, None)

def _reset_slide(i: int):
    original = st.session_state.structure["slides"][i]
    st.session_state.slide_edits.pop(i, None)
    st.session_state.slide_statuses.pop(i, None)
    st.session_state.card_modes[i] = None
    _clear_preview_bytes(i)
    if (original.get("type") == "text"
            and original.get("image", {}).get("needed")
            and original.get("image", {}).get("search_queries")):
        st.session_state.pending_image_fetches.add(i)

def _close_card(i: int):
    st.session_state.card_modes[i] = None

def _clear_preview_bytes(i: int):
    st.session_state.slide_preview_images.pop(i, None)
    try: invalidate_preview_cache(i)
    except Exception: pass


# ════════════════════════════════════════════════════════════════════════════
# SLIDES PIPELINE HELPERS
# ════════════════════════════════════════════════════════════════════════════

def _get_final_slides_for_slides_pipeline() -> list:
    structure = st.session_state.slides_structure or []
    result = []
    for i, slide in enumerate(structure):
        vo  = st.session_state.get(f"sv_{i}", slide.get("voiceover", ""))
        wc  = max(1, len(vo.split()))
        dur = max(6, min(120, int(wc / 2.4)))
        result.append({**slide, "voiceover": vo, "slide_duration_seconds": dur})
    return result

def _clear_slides_session():
    for k in ["slides_pptx_bytes", "slides_pptx_filename", "slides_pptx_path",
              "slides_context", "slides_structure", "slides_edits"]:
        st.session_state.pop(k, None)
    n = len(st.session_state.get("slides_structure") or [])
    for i in range(max(n, 200)):
        st.session_state.pop(f"sv_{i}", None)


# ════════════════════════════════════════════════════════════════════════════
# TOPIC PIPELINE HELPERS
# ════════════════════════════════════════════════════════════════════════════

def _get_topic_scene(i: int) -> dict:
    return st.session_state.topic_scene_edits.get(i, st.session_state.topic_scenes[i])

def _get_final_topic_scenes() -> list:
    return [_get_topic_scene(i) for i in range(len(st.session_state.topic_scenes))]

def _topic_status(i: int) -> str:
    return st.session_state.topic_scene_statuses.get(i, "original")

def _clear_topic_session():
    st.session_state.topic_topic    = ""
    st.session_state.topic_context  = ""
    st.session_state.topic_minutes  = 7
    st.session_state.topic_structure      = None
    st.session_state.topic_scenes         = None
    st.session_state.topic_scene_edits    = {}
    st.session_state.topic_scene_statuses = {}
    st.session_state.topic_card_modes     = {}
    st.session_state.topic_regen_queue    = {}


# ════════════════════════════════════════════════════════════════════════════
# IMAGE SECTION + AUTO-PREFETCH (Full AI pipeline)
# ════════════════════════════════════════════════════════════════════════════

def _render_image_section(i: int, sl: dict):
    has_custom    = i in st.session_state.slide_custom_images
    preview_bytes = st.session_state.slide_preview_images.get(i)
    has_preview   = isinstance(preview_bytes, bytes) and len(preview_bytes) > 100
    image_spec    = sl.get("image", {})
    has_queries   = bool(image_spec.get("search_queries"))

    with st.expander("🖼️ Slide Image", expanded=(has_custom or has_preview)):
        img_col, ctrl_col = st.columns([1, 2])
        with img_col:
            if has_custom:
                st.markdown('<span class="img-badge ib-custom">📁 Custom Upload</span>', unsafe_allow_html=True)
                st.image(st.session_state.slide_custom_images[i], use_container_width=True)
            elif has_preview:
                st.markdown('<span class="img-badge ib-fetched">🔍 API Fetched + Vision ✓</span>', unsafe_allow_html=True)
                st.image(preview_bytes, use_container_width=True)
            else:
                st.markdown(
                    '<div class="img-placeholder"><div style="font-size:24px;margin-bottom:6px;opacity:.4">🖼️</div>'
                    'No image yet.<br>Fetch from API or upload your own.</div>',
                    unsafe_allow_html=True,
                )
        with ctrl_col:
            st.markdown("<div style='font-size:11px;color:#475569;font-weight:600;text-transform:uppercase;letter-spacing:.06em;margin-bottom:10px'>Image Controls</div>", unsafe_allow_html=True)
            if has_queries and not has_custom:
                is_refetch = has_preview
                btn_label  = "🔄 Re-fetch (fresh from Google)" if is_refetch else "🔍 Fetch from API"
                btn_help   = ("Clears ALL caches, fetches a different image from Google + Gemini Vision check."
                              if is_refetch else "Search Google CSE + Gemini Vision filter.")
                if st.button(btn_label, key=f"fetch_{i}", use_container_width=True, help=btn_help):
                    st.session_state.slide_preview_images.pop(i, None)
                    with st.spinner("🔍 Fetching from Google + Vision check…"):
                        ok = _fetch_and_store_preview(i, sl, topic=st.session_state.get("topic", ""), force_refresh=is_refetch)
                    if ok:
                        st.success("✅ New image found and verified!")
                    else:
                        st.warning("No relevant image passed the vision check. Try uploading your own.")
                    st.rerun()
            elif not has_queries:
                st.markdown('<div style="background:#0d1525;border:1px solid #1e3048;border-radius:8px;padding:10px 12px;font-size:11px;color:#3d5470">ℹ️ No search queries defined.<br>Use <b>✏️ Edit</b> to add specific queries, or upload your own.</div>', unsafe_allow_html=True)

            st.markdown("<div style='font-size:11px;color:#3d5070;margin:12px 0 4px 0;font-weight:600'>📁 Upload your own image</div>", unsafe_allow_html=True)
            uploaded = st.file_uploader("Upload image", type=["jpg", "jpeg", "png", "webp"],
                                         key=f"img_up_{i}", label_visibility="collapsed")
            if uploaded is not None:
                raw = uploaded.read()
                if raw:
                    st.session_state.slide_custom_images[i] = raw
                    st.session_state.slide_preview_images.pop(i, None)
                    st.rerun()
            if has_custom:
                st.markdown("<div style='margin-top:8px'></div>", unsafe_allow_html=True)
                if st.button("🗑️ Remove Custom Image", key=f"rmimg_{i}", use_container_width=True):
                    st.session_state.slide_custom_images.pop(i, None)
                    st.rerun()


def _auto_prefetch_previews():
    from concurrent.futures import ThreadPoolExecutor, as_completed as _ac
    slides = st.session_state.structure["slides"]
    topic  = st.session_state.get("topic", "")
    to_fetch = [
        (i, slides[i]) for i in range(len(slides))
        if slides[i].get("type") == "text"
        and slides[i].get("image", {}).get("needed")
        and bool(slides[i].get("image", {}).get("search_queries"))
        and i not in st.session_state.slide_custom_images
        and i not in st.session_state.slide_preview_images
    ]
    if not to_fetch:
        st.session_state.images_prefetched = True
        return
    prog = st.progress(0, text=f"🔍 Auto-fetching images for {len(to_fetch)} slides…")
    done = [0]
    def _fetch_one(args):
        idx, slide = args
        path = fetch_preview_image(idx, slide, topic=topic, force_refresh=False)
        if path and os.path.exists(path) and os.path.getsize(path) > 100:
            return idx, _read_image_bytes(path)
        return idx, None
    results = {}
    with ThreadPoolExecutor(max_workers=4) as ex:
        futures = {ex.submit(_fetch_one, arg): arg[0] for arg in to_fetch}
        for future in _ac(futures):
            idx, img_bytes = future.result()
            if img_bytes: results[idx] = img_bytes
            done[0] += 1
            prog.progress(done[0] / len(to_fetch), text=f"🔍 Images: {done[0]}/{len(to_fetch)} done…")
    for idx, img_bytes in results.items():
        st.session_state.slide_preview_images[idx] = img_bytes
    prog.empty()
    st.session_state.images_prefetched = True
    fetched = len(results); skipped = len(to_fetch) - fetched
    msg = f"✅ Auto-fetched {fetched}/{len(to_fetch)} slide images"
    if skipped: msg += f" — {skipped} had no match (use Re-fetch or upload your own)"
    st.success(msg)


def _process_pending_fetches():
    pending  = st.session_state.pending_image_fetches
    if not pending: return False
    topic    = st.session_state.get("topic", "")
    did_work = False
    for idx in list(pending):
        sl = _get_slide(idx)
        if sl.get("type") != "text" or not sl.get("image", {}).get("search_queries"):
            pending.discard(idx); continue
        with st.spinner(f"🔍 Fetching updated image for slide {idx+1}…"):
            ok = _fetch_and_store_preview(idx, sl, topic=topic, force_refresh=True)
        if ok: st.success(f"✅ Slide {idx+1}: new image ready.")
        else:  st.warning(f"Slide {idx+1}: no matching image — upload your own if needed.")
        pending.discard(idx); did_work = True
    return did_work


# ════════════════════════════════════════════════════════════════════════════
# PHASE 1 — UPLOAD  (all three pipelines)
# ════════════════════════════════════════════════════════════════════════════

def _phase_upload():
    with st.sidebar:
        _show_logo_sidebar()
        st.markdown("## 🎓 Lecture Generator")
        st.markdown("---")

        st.markdown(
            "<div style='font-size:11px;color:#475569;font-weight:700;"
            "text-transform:uppercase;letter-spacing:.06em;margin-bottom:8px'>"
            "Select Pipeline</div>",
            unsafe_allow_html=True,
        )
        pipeline_choice = st.radio(
            "pipeline_selector",
            ["🎬 Full AI Generation", "📊 Slides + Voiceover", "⚡ Topic Video (Pure Manim)"],
            key="up_pipeline_radio",
            label_visibility="collapsed",
        )
        is_slides = "Slides" in pipeline_choice
        is_topic  = "Topic" in pipeline_choice
        st.markdown("---")

        gen_btn = False

        if is_topic:
            # ── Topic Video pipeline inputs ──
            st.markdown('<div class="upload-header">🎯 Topic</div>'
                        '<div class="upload-hint">Required — what should the video teach?</div>',
                        unsafe_allow_html=True)
            st.text_input("topic_input", key="up_topic_topic",
                          placeholder="e.g. Gradient Descent, Newton's Laws, LSTM Networks",
                          label_visibility="collapsed")
            st.markdown("<br>", unsafe_allow_html=True)

            st.markdown('<div class="upload-header">📝 Context / Audience</div>'
                        '<div class="upload-hint">Optional — helps tailor depth and examples.</div>',
                        unsafe_allow_html=True)
            st.text_area("topic_context_area", key="up_topic_context",
                         placeholder="e.g. 2nd-year undergrad ML course. Focus on intuition, include worked examples.",
                         height=80, label_visibility="collapsed")
            st.markdown("<br>", unsafe_allow_html=True)

            st.markdown("**🕐 Target Duration**")
            st.slider("target_minutes", min_value=3, max_value=12, value=7,
                      key="up_topic_minutes", label_visibility="collapsed",
                      help="Gemini will plan scenes to fill this duration.")
            st.caption(f"~{st.session_state.get('up_topic_minutes', 7)} minutes")
            st.markdown("---")
            gen_btn = st.button("⚡ Generate Scene Plan", use_container_width=True,
                                type="primary", key="up_gen_topic")

        elif is_slides:
            # ── Slides pipeline inputs ──
            st.markdown('<div class="upload-header">📎 Upload Your PPTX</div>'
                        '<div class="upload-hint">Required — your slides are used as-is.</div>',
                        unsafe_allow_html=True)
            st.file_uploader("pptx_uploader", type=["pptx"], accept_multiple_files=False,
                             label_visibility="collapsed", key="up_pptx")
            st.markdown("<br>", unsafe_allow_html=True)

            st.markdown('<div class="upload-header">📝 Topic / Context</div>'
                        '<div class="upload-hint">Optional — helps Gemini write better narration.</div>',
                        unsafe_allow_html=True)
            st.text_area("context_area", key="up_context_slides",
                         placeholder="e.g. This lecture covers gradient descent for a 2nd-year ML course.",
                         height=90, label_visibility="collapsed")
            st.markdown("<br>", unsafe_allow_html=True)

            st.markdown("**Optional Reference Sources**")
            for up_key, label, hint, types in [
                ("uf_tr_s", "🎙️ Class Transcript", "Voiceover will echo the teacher's own phrasing.", ["pdf","docx","txt","md"]),
                ("uf_nt_s", "📚 Notes / Textbook", "Adds depth and definitions to the narration.",    ["pdf","docx","txt","md"]),
            ]:
                st.markdown(f'<div class="upload-header">{label}</div>'
                            f'<div class="upload-hint">{hint}</div>', unsafe_allow_html=True)
                st.file_uploader(up_key, type=types, accept_multiple_files=True,
                                 label_visibility="collapsed", key=up_key)
                st.markdown("<br>", unsafe_allow_html=True)

            st.markdown("---")
            gen_btn = st.button("🎙️ Generate Voiceovers", use_container_width=True,
                                type="primary", key="up_gen_slides")
            if st.session_state.get("up_pptx"):
                st.caption(f"📎 {st.session_state.up_pptx.name}")

        else:
            # ── Full AI pipeline inputs ──
            st.text_input("📌 Topic", placeholder="e.g. Robustness in Deep Learning",
                          key="up_topic", help="Leave blank to infer from uploads.")
            st.caption("Optional — inferred from uploads if left blank.")
            st.markdown("---")

            for up_key, label, hint, types in [
                ("uf_tr", "🎙️ Class Transcript", "Teacher's spoken words — voiceover follows what was said.", ["pdf","docx","txt","md"]),
                ("uf_sl", "📊 Slide Deck",        "Class slides — used to structure the video.",              ["pdf","pptx","docx","txt"]),
                ("uf_nt", "📚 Notes / Textbook",  "Adds academic depth and definitions.",                     ["pdf","docx","txt","md"]),
            ]:
                st.markdown(f'<div class="upload-header">{label}</div>'
                            f'<div class="upload-hint">{hint}</div>', unsafe_allow_html=True)
                st.file_uploader(up_key, type=types, accept_multiple_files=True,
                                 label_visibility="collapsed", key=up_key)
                st.markdown("<br>", unsafe_allow_html=True)

            st.toggle("⚡ Skip Manim animations", value=False, key="up_skip")
            st.markdown("---")
            gen_btn = st.button("🔍 Generate Structure", use_container_width=True,
                                type="primary", key="up_gen_full")

            tr_files = st.session_state.get("uf_tr", []) or []
            sl_files = st.session_state.get("uf_sl", []) or []
            nt_files = st.session_state.get("uf_nt", []) or []
            total = len(tr_files) + len(sl_files) + len(nt_files)
            if total:
                st.markdown(f"**{total} file(s) loaded**")
                for f in tr_files: st.caption(f"🎙️ {f.name}")
                for f in sl_files: st.caption(f"📊 {f.name}")
                for f in nt_files: st.caption(f"📚 {f.name}")

    # ── Main area ─────────────────────────────────────────────────────────
    _show_logo_topright()
    st.title("🎓 AI Lecture & Video Generator")

    if is_topic:
        st.caption("Enter a topic → AI plans Manim scenes → Review & edit → Render pure animated MP4.")
        _render_upload_main_topic()
    elif is_slides:
        st.caption("Upload your existing PPTX → AI generates voiceover narration → Download narrated MP4.")
        _render_upload_main_slides()
    else:
        st.caption("Upload class material → AI generates structure → Review & edit slides → Render narrated MP4.")
        _render_upload_main_full()

    if gen_btn:
        if is_topic:
            _handle_generate_topic()
        elif is_slides:
            _handle_generate_slides()
        else:
            _handle_generate_full()


def _render_upload_main_topic():
    c1, c2, c3, c4 = st.columns(4)
    for col, icon, lbl, detail in zip(
        [c1,c2,c3,c4],
        ["⚡","🎨","🎧","🎬"],
        ["AI Scene Plan","Pure Manim","Edge TTS","MP4 Output"],
        ["Gemini designs each scene","23 animation templates","Neural voice narration","Full HD 1920×1080"],
    ):
        with col:
            st.markdown(f'<div class="metric-card"><div style="font-size:1.6rem;margin-bottom:4px">{icon}</div>'
                        f'<div class="metric-value" style="font-size:0.88rem">{lbl}</div>'
                        f'<div class="metric-label">{detail}</div></div>', unsafe_allow_html=True)
    st.markdown('<div class="section-sep"></div>', unsafe_allow_html=True)
    with st.expander("📖 How the Topic Video pipeline works", expanded=True):
        st.markdown("""
**100% Manim animated — no PowerPoint, no photos.**

| Step | What happens |
|------|-------------|
| 🎯 Enter topic | Any subject — math, physics, CS, biology… |
| 🤖 Gemini plans | AI chooses the best animation template for each concept |
| ✏️ Review scenes | Edit voiceovers, swap scene types, tweak params |
| 🎨 Manim renders | Each scene animated with 23 different templates |
| 🎧 Edge TTS | Neural voiceover synced to each animation |
| 🎬 Compose | All scenes concatenated into a final MP4 |

**Available scene types include:** `formula`, `worked_example`, `calculus_plot`,
`neural_network`, `block_flow`, `concept_map`, `bar_chart`, `timeline`,
`step_by_step`, `definition_box`, `comparison_table`, `physics_forces`, `gradient_descent`, and more.

**Perfect for:** math derivations, algorithm explanations, physics problems, CS concepts.
        """)
    st.info("👈 Enter a topic in the sidebar, then click **⚡ Generate Scene Plan**.", icon="💡")


def _render_upload_main_slides():
    c1, c2, c3, c4 = st.columns(4)
    for col, icon, lbl, detail in zip(
        [c1,c2,c3,c4],
        ["📎","🤖","🎧","🎬"],
        ["Your PPTX","AI Narration","Edge TTS","MP4 Output"],
        ["Slides rendered as-is","Gemini scripts per slide","Neural voice narration","Full HD 1920×1080"],
    ):
        with col:
            st.markdown(f'<div class="metric-card"><div style="font-size:1.6rem;margin-bottom:4px">{icon}</div>'
                        f'<div class="metric-value" style="font-size:0.88rem">{lbl}</div>'
                        f'<div class="metric-label">{detail}</div></div>', unsafe_allow_html=True)
    st.markdown('<div class="section-sep"></div>', unsafe_allow_html=True)
    with st.expander("📖 How the Slides + Voiceover pipeline works", expanded=True):
        st.markdown("""
**This pipeline leaves your slides completely unchanged.**

| Step | What happens |
|------|-------------|
| 📎 Upload PPTX | Your presentation is preserved exactly |
| 🤖 Gemini narration | AI writes a voiceover script for each slide |
| ✏️ Review & edit | You can tweak any narration before rendering |
| 🎧 Edge TTS | Each voiceover is converted to speech |
| 🎬 Compose | Slide images + audio are merged into an MP4 |
        """)
    st.info("👈 Upload your PPTX in the sidebar, then click **🎙️ Generate Voiceovers**.", icon="💡")


def _render_upload_main_full():
    c1, c2, c3, c4 = st.columns(4)
    for col, icon, lbl, detail in zip(
        [c1,c2,c3,c4],
        ["🤖","🎬","🖼️","🎧"],
        ["AI Script","Manim Diagrams","Vision-Filtered Images","TTS Narration"],
        ["From your uploads","23 data-driven templates","Google CSE + Gemini Vision","Parallel Edge TTS, synced"],
    ):
        with col:
            st.markdown(f'<div class="metric-card"><div style="font-size:1.6rem;margin-bottom:4px">{icon}</div>'
                        f'<div class="metric-value" style="font-size:0.88rem">{lbl}</div>'
                        f'<div class="metric-label">{detail}</div></div>', unsafe_allow_html=True)
    st.markdown('<div class="section-sep"></div>', unsafe_allow_html=True)
    with st.expander("📖 How it works", expanded=True):
        st.markdown("""
**Three upload categories — all optional:**

| Upload | Purpose |
|--------|---------|
| 🎙️ **Transcript** | Voiceover uses teacher's own phrasing |
| 📊 **Slide Deck** | Sets slide structure and topic order |
| 📚 **Notes / Textbook** | Adds depth, definitions, formal language |

**Workflow:** 🔍 Generate → 🖼️ Auto-fetch images → ✏️ Review & edit → 🚀 Render MP4
        """)
    st.info("👈 Upload files in the sidebar, then click **Generate Structure**.", icon="💡")


# ── Generate handlers ──────────────────────────────────────────────────────

def _handle_generate_topic():
    topic = (st.session_state.get("up_topic_topic") or "").strip()
    if not topic:
        st.error("⚠️ Please enter a topic first.")
        return
    context = (st.session_state.get("up_topic_context") or "").strip()
    minutes = int(st.session_state.get("up_topic_minutes") or 7)

    with st.spinner(f"🤖 Gemini is designing a {minutes}-minute scene plan for '{topic}'…"):
        try:
            structure = generate_topic_structure(topic=topic, context=context, target_minutes=minutes)
        except Exception as e:
            st.error(f"Scene plan generation failed: {e}")
            with st.expander("Details"): st.code(traceback.format_exc())
            return

    scenes = structure.get("scenes", [])
    if not scenes:
        st.error("Gemini returned an empty scene plan. Try again or rephrase the topic.")
        return

    n        = len(scenes)
    est_mins = structure.get("estimated_minutes", minutes)
    st.success(f"✅ Scene plan ready — {n} scenes (~{est_mins:.1f} min)")

    st.session_state.update(
        pipeline              = "topic",
        topic_topic           = topic,
        topic_context         = context,
        topic_minutes         = minutes,
        topic_structure       = structure,
        topic_scenes          = scenes,
        topic_scene_edits     = {},
        topic_scene_statuses  = {i: "original" for i in range(n)},
        topic_card_modes      = {i: None for i in range(n)},
        topic_regen_queue     = {},
        review_phase          = "topic_review",
        phase                 = "topic_review",
    )
    st.rerun()


def _handle_generate_slides():
    pptx_file = st.session_state.get("up_pptx")
    if not pptx_file:
        st.error("⚠️ Please upload a PPTX file first.")
        return
    pptx_bytes = pptx_file.read()
    if not pptx_bytes:
        st.error("⚠️ The uploaded PPTX appears to be empty.")
        return

    os.makedirs(OUTPUT_DIR, exist_ok=True)
    pptx_save_path = os.path.join(OUTPUT_DIR, "uploaded_slides.pptx")
    with open(pptx_save_path, "wb") as f:
        f.write(pptx_bytes)

    tr_files = st.session_state.get("uf_tr_s", []) or []
    nt_files = st.session_state.get("uf_nt_s", []) or []
    tr_text = nt_text = None
    with st.spinner("Reading reference files…"):
        if tr_files:
            t = _extract(tr_files, "transcript"); tr_text = t or None
            if t: st.success(f"🎙️ Transcript: {len(t):,} chars")
        if nt_files:
            t = _extract(nt_files, "notes"); nt_text = t or None
            if t: st.success(f"📚 Notes: {len(t):,} chars")

    context_text = (st.session_state.get("up_context_slides") or "").strip()

    with st.spinner(f"🤖 Extracting slides and generating voiceovers…"):
        try:
            titles = extract_slide_titles_from_pptx(pptx_save_path)
            slides = generate_voiceovers_for_slide_titles(titles, tr_text, nt_text, context_text or None)
        except Exception as e:
            st.error(f"Generation failed: {e}")
            with st.expander("Details"): st.code(traceback.format_exc())
            return

    for i, s in enumerate(slides):
        st.session_state[f"sv_{i}"] = s.get("voiceover", "")

    n       = len(slides)
    est_min = max(1, sum(s.get("slide_duration_seconds", 8) for s in slides) // 60)
    st.success(f"✅ Generated voiceovers for {n} slides (~{est_min} min)")

    st.session_state.update(
        pipeline             = "slides",
        slides_pptx_bytes    = pptx_bytes,
        slides_pptx_filename = pptx_file.name,
        slides_pptx_path     = pptx_save_path,
        slides_context       = context_text,
        slides_structure     = slides,
        slides_edits         = {},
        transcript_text      = tr_text,
        notes_text           = nt_text,
        review_phase         = "slides_review",
        phase                = "slides_review",
    )
    st.rerun()


def _handle_generate_full():
    topic    = (st.session_state.get("up_topic") or "").strip()
    tr_files = st.session_state.get("uf_tr", []) or []
    sl_files = st.session_state.get("uf_sl", []) or []
    nt_files = st.session_state.get("uf_nt", []) or []
    skip     = st.session_state.get("up_skip", False)

    if not topic and not any([tr_files, sl_files, nt_files]):
        st.error("⚠️ Please enter a topic or upload at least one file.")
        return

    tr_text = sl_text = nt_text = None
    with st.spinner("Reading uploaded files..."):
        if tr_files:
            t = _extract(tr_files, "transcript"); tr_text = t or None
            if t: st.success(f"🎙️ Transcript: {len(t):,} chars")
        if sl_files:
            t = _extract(sl_files, "slides"); sl_text = t or None
            if t: st.success(f"📊 Slides: {len(t):,} chars")
        if nt_files:
            t = _extract(nt_files, "notes"); nt_text = t or None
            if t: st.success(f"📚 Notes: {len(t):,} chars")

    if not topic and not any([tr_text, sl_text, nt_text]):
        st.error("No readable content found. Check your files or enter a topic.")
        return

    with st.spinner("🤖 Calling Gemini to generate lecture structure…"):
        try:
            structure = generate_lecture_structure(topic=topic, transcript_text=tr_text,
                                                    slides_text=sl_text, notes_text=nt_text)
        except Exception as e:
            st.error(f"Structure generation failed: {e}")
            with st.expander("Details"): st.code(traceback.format_exc())
            return

    n = len(structure["slides"])
    st.session_state.update(
        pipeline              = "full",
        structure             = structure,
        topic                 = topic,
        transcript_text       = tr_text,
        slides_text           = sl_text,
        notes_text            = nt_text,
        skip_diagrams         = skip,
        slide_statuses        = {i: "original" for i in range(n)},
        slide_edits           = {},
        card_modes            = {i: None for i in range(n)},
        regen_queue           = {},
        slide_preview_images  = {},
        slide_custom_images   = {},
        images_prefetched     = False,
        pending_image_fetches = set(),
        review_phase          = "review",
        phase                 = "review",
    )
    st.rerun()


# ════════════════════════════════════════════════════════════════════════════
# PHASE 2A — FULL AI PIPELINE REVIEW
# ════════════════════════════════════════════════════════════════════════════

def _render_card(i: int):
    sl     = _get_slide(i)
    status = _status(i)
    mode   = st.session_state.card_modes.get(i)

    border  = {"original":"#243048","edited":"#0369a1","regenerated":"#7c3aed"}.get(status,"#243048")
    stype   = sl.get("type","text")
    dur     = sl.get("slide_duration_seconds", 90)
    title   = sl.get("title", f"Slide {i+1}")
    bullets = sl.get("bullets", [])

    status_badge = {"original":'<span class="s-badge b-orig">⚪ Original</span>',
                    "edited":'<span class="s-badge b-edited">✏️ Edited</span>',
                    "regenerated":'<span class="s-badge b-regen">🔄 Regen</span>'}.get(status,"")
    type_badge = ('<span class="s-badge b-diagram">📊 Diagram</span>' if stype=="diagram"
                  else '<span class="s-badge b-text">📝 Text</span>')
    img_indicator = ""
    if stype == "text":
        if i in st.session_state.slide_custom_images:
            img_indicator = '<span class="s-badge ib-custom" style="font-size:9px;padding:2px 7px">📁 Custom</span>'
        elif i in st.session_state.slide_preview_images:
            img_indicator = '<span class="s-badge ib-fetched" style="font-size:9px;padding:2px 7px">🖼️ Image</span>'

    st.markdown(
        f'<div class="scard" style="border-left:3px solid {border}">'
        f'  <span class="s-num">SLIDE {i+1:02d}</span>&nbsp;'
        f'  <span class="s-title">{title[:68]}</span>'
        f'  {type_badge}{status_badge}{img_indicator}'
        f'  <span class="b-dur">⏱ {dur}s</span></div>',
        unsafe_allow_html=True,
    )
    for b in bullets[:2]:
        st.markdown(f'<div class="bullet-prev">▸ {b[:100]}</div>', unsafe_allow_html=True)
    if len(bullets) > 2:
        st.markdown(f'<div class="bullet-prev" style="color:#2a3a52">+{len(bullets)-2} more…</div>', unsafe_allow_html=True)

    if stype == "text" and mode is None:
        _render_image_section(i, sl)

    if mode is None:
        ca, cb, cc, cd, _ = st.columns([1, 1, 1, 1, 5])
        with ca:
            if st.button("✅ Keep", key=f"keep_{i}"): st.session_state.card_modes[i]=None; st.rerun()
        with cb:
            if st.button("✏️ Edit", key=f"edit_btn_{i}"): _open_edit(i); st.rerun()
        with cc:
            if st.button("🔄 Regen", key=f"regen_btn_{i}"): st.session_state.card_modes[i]="regen"; st.rerun()
        with cd:
            if status != "original":
                if st.button("↩️ Reset", key=f"reset_{i}"): _reset_slide(i); st.rerun()

    elif mode == "edit":
        st.markdown('<div class="panel"><div class="panel-title">✏️ Edit Slide Content</div></div>', unsafe_allow_html=True)
        st.text_input("Title  (max 55 chars)", key=f"et_{i}", max_chars=55)
        col_l, col_r = st.columns([3, 2])
        with col_l: st.text_area("Bullets  (one per line)", key=f"eb_{i}", height=130)
        with col_r: st.text_area("Voiceover / Narration  (120–200 words)", key=f"ev_{i}", height=130)
        tc1, tc2 = st.columns(2)
        with tc1: stype_sel = st.selectbox("Slide type", ["text","diagram"], key=f"ety_{i}")
        with tc2:
            diag_opts = ["none"] + sorted(AVAILABLE_DIAGRAM_TYPES)
            st.selectbox("Diagram type", diag_opts, key=f"ed_{i}", disabled=(stype_sel=="text"))
        if stype_sel == "text":
            img_needed = st.checkbox("🖼️ Include image for this slide", key=f"ein_{i}")
            if img_needed:
                st.text_area("Image search queries  (one per line — be very specific)", key=f"eiq_{i}", height=65)
                eu = st.file_uploader("Edit upload", type=["jpg","jpeg","png","webp"],
                                       key=f"img_up_edit_{i}", label_visibility="collapsed")
                if eu is not None:
                    raw = eu.read()
                    if raw:
                        st.session_state.slide_custom_images[i] = raw
                        st.session_state.slide_preview_images.pop(i, None)
        else:
            st.info("ℹ️ Diagram slides use animated Manim visuals instead of photos.")
        bc1, bc2, bc3, _ = st.columns([1, 1, 1, 5])
        with bc1:
            if st.button("💾 Save", key=f"save_{i}", type="primary"): _save_edit(i); st.rerun()
        with bc2:
            if st.button("↩️ Reset", key=f"reset_edit_{i}"): _reset_slide(i); st.rerun()
        with bc3:
            if st.button("❌ Cancel", key=f"cancel_edit_{i}"): _close_card(i); st.rerun()

    elif mode == "regen":
        if i in st.session_state.regen_queue:
            rp = st.session_state.regen_queue[i]
            with st.spinner(f"🤖 Regenerating slide {i+1}…"):
                try:
                    slides_all  = st.session_state.structure["slides"]
                    surrounding = [slides_all[j] for j in range(max(0,i-2),min(len(slides_all),i+3)) if j!=i]
                    result = regenerate_single_slide(
                        slide_idx=i, current_slide=_get_slide(i),
                        structure_metadata=st.session_state.structure["metadata"],
                        surrounding_slides=surrounding,
                        transcript_text=st.session_state.transcript_text,
                        slides_text=st.session_state.slides_text,
                        notes_text=st.session_state.notes_text,
                        mode=rp["mode"], custom_prompt=rp["custom_prompt"],
                    )
                    st.session_state.slide_edits[i]    = result
                    st.session_state.slide_statuses[i] = "regenerated"
                    st.session_state.card_modes[i]     = None
                    del st.session_state.regen_queue[i]
                    _clear_preview_bytes(i)
                    if (result.get("type")=="text" and result.get("image",{}).get("needed")
                            and result.get("image",{}).get("search_queries")):
                        st.session_state.pending_image_fetches.add(i)
                    st.success(f"✅ Slide {i+1} regenerated!")
                except Exception as e:
                    st.error(f"Regeneration failed: {e}")
                    with st.expander("Details"): st.code(traceback.format_exc())
                    del st.session_state.regen_queue[i]
                    st.session_state.card_modes[i] = None
            st.rerun(); return

        st.markdown('<div class="panel"><div class="panel-title">🔄 Regenerate with AI</div></div>', unsafe_allow_html=True)
        custom_prompt = st.text_area("📝 Additional instructions  (optional)", key=f"rp_{i}",
                                      placeholder="e.g. Add a real-world example, use simpler language…", height=90)
        regen_mode = st.radio("What to regenerate",
                               ["🔁 Full slide  (title + bullets + voiceover)","🎙️ Voiceover only"],
                               key=f"rm_{i}", horizontal=True)
        mode_val = "full" if "Full" in regen_mode else "voiceover_only"
        if mode_val == "voiceover_only": st.caption("ℹ️ Title and bullets kept as-is.")
        rc1, rc2, _ = st.columns([1.2, 1, 6])
        with rc1:
            if st.button("🚀 Regenerate", key=f"do_regen_{i}", type="primary"):
                st.session_state.regen_queue[i] = {"mode":mode_val,"custom_prompt":custom_prompt}; st.rerun()
        with rc2:
            if st.button("❌ Cancel", key=f"cancel_regen_{i}"): _close_card(i); st.rerun()

    st.markdown("<div style='height:6px'></div>", unsafe_allow_html=True)


def _phase_review():
    slides = st.session_state.structure["slides"]
    meta   = st.session_state.structure.get("metadata", {})
    n      = len(slides)

    if not st.session_state.get("images_prefetched", False):
        _auto_prefetch_previews(); st.rerun()

    if st.session_state.pending_image_fetches:
        if _process_pending_fetches(): st.rerun()

    n_mod  = _n_modified()
    n_reg  = sum(1 for v in st.session_state.slide_statuses.values() if v=="regenerated")
    n_ed   = sum(1 for v in st.session_state.slide_statuses.values() if v=="edited")
    n_cust = len(st.session_state.slide_custom_images)
    n_img  = len(st.session_state.slide_preview_images)

    with st.sidebar:
        _show_logo_sidebar()
        st.markdown("## 🎓 Lecture Generator")
        st.markdown("---")
        st.markdown(f"**📊 {n} slides** &nbsp;·&nbsp; ⏱ ~{meta.get('estimated_duration_minutes','?')} min")
        st.markdown(f"👤 {str(meta.get('target_audience','?')).title()} level")
        if st.session_state.topic: st.markdown(f"📌 _{st.session_state.topic}_")
        st.markdown("---")
        st.markdown(f'<div style="font-size:12px;color:#475569;margin-bottom:10px">✏️ {n_ed} edited &nbsp;🔄 {n_reg} regen &nbsp;📁 {n_cust} custom &nbsp;🖼️ {n_img} imgs</div>', unsafe_allow_html=True)
        st.session_state.skip_diagrams = st.toggle("⚡ Skip Manim", value=st.session_state.skip_diagrams, key="rv_skip")
        st.markdown("---")
        if st.button("🚀 Finalize & Render", use_container_width=True, type="primary", key="fin_sidebar"):
            st.session_state.phase = "rendering"; st.rerun()
        st.markdown("<br>", unsafe_allow_html=True)
        if st.button("🔄 Start Over", use_container_width=True, key="start_over"):
            for k in ["structure","slide_statuses","slide_edits","card_modes","regen_queue",
                      "slide_preview_images","slide_custom_images","images_prefetched","pending_image_fetches"]:
                st.session_state.pop(k, None)
            st.session_state.phase = "upload"; st.rerun()

    _show_logo_topright()
    st.title("📋 Review & Edit Slides")
    mc1, mc2, mc3, mc4 = st.columns(4)
    with mc1: st.markdown(f'<div class="metric-card"><div class="metric-value">{n}</div><div class="metric-label">Total Slides</div></div>', unsafe_allow_html=True)
    with mc2: st.markdown(f'<div class="metric-card"><div class="metric-value">{n_mod}</div><div class="metric-label">Modified</div></div>', unsafe_allow_html=True)
    with mc3: st.markdown(f'<div class="metric-card"><div class="metric-value" style="font-size:1rem"><span style="color:#60a5fa">{n_img}</span> 🖼️ &nbsp;<span style="color:#4ade80">{n_cust}</span> 📁</div><div class="metric-label">Images / Custom</div></div>', unsafe_allow_html=True)
    with mc4: fin_top = st.button("🚀 Finalize & Render Video", type="primary", use_container_width=True, key="fin_top")

    if st.session_state.topic:
        est = meta.get("estimated_duration_minutes","?"); aud = str(meta.get("target_audience","?")).title()
        st.caption(f"📌 **{st.session_state.topic}** · {aud} · ~{est} min")

    st.markdown('<div class="section-sep"></div>', unsafe_allow_html=True)
    st.markdown('<div class="phase-banner">Images auto-fetch on load. <b>🔄 Re-fetch</b> gets a fresh result from Google. <b>✏️ Edit</b> content &amp; queries. <b>🔄 Regen</b> with AI. <b>↩️ Reset</b> to revert.</div>', unsafe_allow_html=True)

    for i in range(n): _render_card(i)

    st.markdown('<div class="section-sep"></div>', unsafe_allow_html=True)
    fin_bot = st.button("🚀 Finalize & Render Video", type="primary", use_container_width=True, key="fin_bot")
    if fin_top or fin_bot:
        st.session_state.phase = "rendering"; st.rerun()


# ════════════════════════════════════════════════════════════════════════════
# PHASE 2B — SLIDES PIPELINE REVIEW
# ════════════════════════════════════════════════════════════════════════════

def _phase_slides_review():
    structure = st.session_state.slides_structure or []
    n         = len(structure)

    for i, slide in enumerate(structure):
        key = f"sv_{i}"
        if key not in st.session_state:
            st.session_state[key] = slide.get("voiceover", "")

    with st.sidebar:
        _show_logo_sidebar()
        st.markdown("## 🎓 Lecture Generator")
        st.markdown("---")
        st.markdown(f"**📎 {n} slides**")
        if st.session_state.slides_pptx_filename:
            st.caption(f"📌 {st.session_state.slides_pptx_filename}")
        total_dur = sum(max(6,min(120,int(max(1,len(st.session_state.get(f"sv_{i}","").split()))/2.4))) for i in range(n))
        st.markdown(f"⏱ ~{max(1, total_dur // 60)} min estimated")
        st.markdown("---")
        if st.button("🚀 Finalize & Render", use_container_width=True, type="primary", key="slides_fin_sidebar"):
            st.session_state.phase = "slides_rendering"; st.rerun()
        st.markdown("<br>", unsafe_allow_html=True)
        if st.button("🔄 Start Over", use_container_width=True, key="slides_start_over"):
            _clear_slides_session(); st.session_state.phase = "upload"; st.rerun()

    _show_logo_topright()
    st.title("📋 Review Slide Voiceovers")
    if st.session_state.slides_pptx_filename:
        st.caption(f"📎 **{st.session_state.slides_pptx_filename}** · {n} slides · slides unchanged")

    total_dur = sum(max(6,min(120,int(max(1,len(st.session_state.get(f"sv_{i}","").split()))/2.4))) for i in range(n))
    mc1, mc2, mc3 = st.columns(3)
    with mc1: st.markdown(f'<div class="metric-card"><div class="metric-value">{n}</div><div class="metric-label">Slides</div></div>', unsafe_allow_html=True)
    with mc2: st.markdown(f'<div class="metric-card"><div class="metric-value">~{max(1, total_dur // 60)} min</div><div class="metric-label">Est. Duration</div></div>', unsafe_allow_html=True)
    with mc3: fin_top = st.button("🚀 Finalize & Render Video", type="primary", use_container_width=True, key="slides_fin_top")

    st.markdown('<div class="section-sep"></div>', unsafe_allow_html=True)
    st.markdown('<div class="phase-banner">Review and edit the AI-generated narration for each slide. Word count drives clip duration — longer text = more time on each slide. Your <b>original slides are not changed</b>.</div>', unsafe_allow_html=True)

    for i, slide in enumerate(structure):
        title   = slide.get("title", f"Slide {i+1}")
        cur_vo  = st.session_state.get(f"sv_{i}", slide.get("voiceover",""))
        wc      = max(1, len(cur_vo.split()))
        est_dur = max(6, min(120, int(wc / 2.4)))
        with st.expander(f"Slide {i+1:02d} · {title[:65]}   ⏱ ~{est_dur}s", expanded=False):
            st.markdown(f'<div style="font-size:11px;color:#3d5470;margin-bottom:6px"><b>Title:</b> {title}</div>', unsafe_allow_html=True)
            st.text_area("Voiceover / Narration", key=f"sv_{i}", height=130,
                         help="Edit the narration. Word count determines clip duration (~2.4 words/second).")
            new_wc  = max(1, len(st.session_state.get(f"sv_{i}", "").split()))
            new_dur = max(6, min(120, int(new_wc / 2.4)))
            st.caption(f"~{new_dur}s · {new_wc} words")

    st.markdown('<div class="section-sep"></div>', unsafe_allow_html=True)
    fin_bot = st.button("🚀 Finalize & Render Video", type="primary", use_container_width=True, key="slides_fin_bot")
    if fin_top or fin_bot:
        st.session_state.phase = "slides_rendering"; st.rerun()


# ════════════════════════════════════════════════════════════════════════════
# PHASE 2C — TOPIC VIDEO REVIEW
# ════════════════════════════════════════════════════════════════════════════

def _render_topic_scene_card(i: int):
    scene  = _get_topic_scene(i)
    status = _topic_status(i)
    mode   = st.session_state.topic_card_modes.get(i)

    border      = {"original":"#1e3a2a","edited":"#0369a1","regenerated":"#7c3aed"}.get(status,"#1e3a2a")
    stype       = scene.get("scene_type","bullet_list")
    dur         = scene.get("duration_seconds", 30)
    title       = scene.get("title", f"Scene {i+1}")
    voiceover   = scene.get("voiceover","")
    wc          = max(1, len(voiceover.split()))

    status_badge = {"original":'<span class="s-badge b-orig">⚪ Original</span>',
                    "edited":'<span class="s-badge b-edited">✏️ Edited</span>',
                    "regenerated":'<span class="s-badge b-regen">🔄 Regen</span>'}.get(status,"")

    st.markdown(
        f'<div class="scene-card" style="border-left:3px solid {border}">'
        f'  <span class="s-num">SCENE {i+1:02d}</span>&nbsp;'
        f'  <span class="s-title">{title[:68]}</span>'
        f'  <span class="scene-type-badge">🎨 {stype}</span>{status_badge}'
        f'  <span class="b-dur">⏱ {dur}s · {wc}w</span></div>',
        unsafe_allow_html=True,
    )

    # Show voiceover preview
    if voiceover:
        preview = voiceover[:120] + ("…" if len(voiceover) > 120 else "")
        st.markdown(f'<div class="bullet-prev">🎤 {preview}</div>', unsafe_allow_html=True)

    if mode is None:
        ca, cb, cc, _ = st.columns([1, 1, 1, 6])
        with ca:
            if st.button("✅ Keep", key=f"tkeep_{i}"): st.session_state.topic_card_modes[i]=None; st.rerun()
        with cb:
            if st.button("✏️ Edit", key=f"tedit_{i}"):
                st.session_state.topic_card_modes[i] = "edit"
                st.session_state[f"tev_{i}"]  = scene.get("voiceover","")
                st.session_state[f"tet_{i}"]  = title
                st.session_state[f"testy_{i}"] = stype
                st.session_state[f"tdur_{i}"] = dur
                st.rerun()
        with cc:
            if st.button("🔄 Regen", key=f"tregen_{i}"): st.session_state.topic_card_modes[i]="regen"; st.rerun()

    elif mode == "edit":
        st.markdown('<div class="panel"><div class="panel-title">✏️ Edit Scene</div></div>', unsafe_allow_html=True)
        col_l, col_r = st.columns([2, 1])
        with col_l:
            st.text_input("Scene Title", key=f"tet_{i}", max_chars=55)
            st.text_area("Voiceover / Narration  (120–180 words for content scenes)", key=f"tev_{i}", height=150)
        with col_r:
            scene_type_opts = sorted(AVAILABLE_SCENE_TYPES)
            cur_type_idx = scene_type_opts.index(stype) if stype in scene_type_opts else 0
            new_stype = st.selectbox("Scene Type", scene_type_opts, index=cur_type_idx, key=f"testy_{i}")
            st.number_input("Duration (seconds)", min_value=8, max_value=120,
                             value=dur, key=f"tdur_{i}")
            st.caption("Duration auto-adjusts to voiceover length during render.")

        bc1, bc2, bc3, _ = st.columns([1, 1, 1, 5])
        with bc1:
            if st.button("💾 Save", key=f"tsave_{i}", type="primary"):
                new_vo    = st.session_state.get(f"tev_{i}", scene.get("voiceover",""))
                new_title = st.session_state.get(f"tet_{i}", title)[:55]
                new_type  = st.session_state.get(f"testy_{i}", stype)
                new_dur   = int(st.session_state.get(f"tdur_{i}", dur))
                edited    = copy.deepcopy(scene)
                edited["title"]          = new_title
                edited["voiceover"]      = new_vo
                edited["scene_type"]     = new_type
                edited["duration_seconds"] = new_dur
                if "params" in edited:
                    edited["params"]["title"] = new_title
                    edited["params"]["type"]  = new_type
                st.session_state.topic_scene_edits[i]    = edited
                st.session_state.topic_scene_statuses[i] = "edited"
                st.session_state.topic_card_modes[i]     = None
                st.rerun()
        with bc2:
            if st.button("↩️ Reset", key=f"treset_{i}"):
                st.session_state.topic_scene_edits.pop(i, None)
                st.session_state.topic_scene_statuses.pop(i, None)
                st.session_state.topic_card_modes[i] = None; st.rerun()
        with bc3:
            if st.button("❌ Cancel", key=f"tcancel_{i}"):
                st.session_state.topic_card_modes[i] = None; st.rerun()

    elif mode == "regen":
        if i in st.session_state.topic_regen_queue:
            rq = st.session_state.topic_regen_queue[i]
            with st.spinner(f"🤖 Regenerating scene {i+1}…"):
                try:
                    all_scenes  = st.session_state.topic_scenes
                    surrounding = [all_scenes[j] for j in range(max(0,i-2),min(len(all_scenes),i+3)) if j!=i]
                    result = regenerate_topic_scene(
                        scene_idx=i, current_scene=_get_topic_scene(i),
                        topic=st.session_state.topic_topic,
                        context=st.session_state.topic_context,
                        surrounding_scenes=surrounding,
                        custom_prompt=rq["custom_prompt"],
                    )
                    st.session_state.topic_scene_edits[i]    = result
                    st.session_state.topic_scene_statuses[i] = "regenerated"
                    st.session_state.topic_card_modes[i]     = None
                    del st.session_state.topic_regen_queue[i]
                    st.success(f"✅ Scene {i+1} regenerated!")
                except Exception as e:
                    st.error(f"Regeneration failed: {e}")
                    with st.expander("Details"): st.code(traceback.format_exc())
                    del st.session_state.topic_regen_queue[i]
                    st.session_state.topic_card_modes[i] = None
            st.rerun(); return

        st.markdown('<div class="panel"><div class="panel-title">🔄 Regenerate Scene with AI</div></div>', unsafe_allow_html=True)
        custom_prompt = st.text_area("📝 Additional instructions  (optional)", key=f"trp_{i}",
                                      placeholder="e.g. Make it a worked example with step-by-step math, use calculus_plot instead…", height=90)
        rc1, rc2, _ = st.columns([1.2, 1, 6])
        with rc1:
            if st.button("🚀 Regenerate", key=f"tdo_regen_{i}", type="primary"):
                st.session_state.topic_regen_queue[i] = {"custom_prompt": custom_prompt}; st.rerun()
        with rc2:
            if st.button("❌ Cancel", key=f"tcancel_regen_{i}"):
                st.session_state.topic_card_modes[i] = None; st.rerun()

    st.markdown("<div style='height:6px'></div>", unsafe_allow_html=True)


def _phase_topic_review():
    scenes = st.session_state.topic_scenes or []
    n      = len(scenes)
    topic  = st.session_state.topic_topic
    struct = st.session_state.topic_structure or {}

    n_ed  = sum(1 for v in st.session_state.topic_scene_statuses.values() if v=="edited")
    n_reg = sum(1 for v in st.session_state.topic_scene_statuses.values() if v=="regenerated")
    total_dur = sum(_get_topic_scene(i).get("duration_seconds", 30) for i in range(n))
    est_min   = max(1, total_dur // 60)

    with st.sidebar:
        _show_logo_sidebar()
        st.markdown("## 🎓 Lecture Generator")
        st.markdown("---")
        st.markdown(f'<span class="pipeline-pill pp-topic">⚡ Topic Video</span>', unsafe_allow_html=True)
        st.markdown(f"**{n} scenes** · ~{est_min} min")
        st.markdown(f"📌 _{topic}_")
        st.markdown("---")
        st.markdown(f'<div style="font-size:12px;color:#475569;margin-bottom:10px">✏️ {n_ed} edited &nbsp;🔄 {n_reg} regen</div>', unsafe_allow_html=True)
        st.markdown("---")
        if st.button("🚀 Render Topic Video", use_container_width=True, type="primary", key="topic_fin_sidebar"):
            st.session_state.phase = "topic_rendering"; st.rerun()
        st.markdown("<br>", unsafe_allow_html=True)
        if st.button("🔄 Start Over", use_container_width=True, key="topic_start_over"):
            _clear_topic_session(); st.session_state.phase = "upload"; st.rerun()

    _show_logo_topright()
    st.title("🎨 Review Scene Plan")
    st.caption(f"⚡ **{topic}** · Pure Manim · {n} scenes · ~{est_min} min")

    if struct.get("description"):
        st.markdown(f'<div class="phase-banner">📝 {struct["description"]}</div>', unsafe_allow_html=True)

    mc1, mc2, mc3, mc4 = st.columns(4)
    with mc1: st.markdown(f'<div class="metric-card"><div class="metric-value">{n}</div><div class="metric-label">Scenes</div></div>', unsafe_allow_html=True)
    with mc2: st.markdown(f'<div class="metric-card"><div class="metric-value">~{est_min} min</div><div class="metric-label">Est. Duration</div></div>', unsafe_allow_html=True)
    with mc3: st.markdown(f'<div class="metric-card"><div class="metric-value">{n_ed+n_reg}</div><div class="metric-label">Modified</div></div>', unsafe_allow_html=True)
    with mc4: fin_top = st.button("🚀 Render Topic Video", type="primary", use_container_width=True, key="topic_fin_top")

    st.markdown('<div class="section-sep"></div>', unsafe_allow_html=True)
    st.markdown(
        '<div class="phase-banner">'
        "Each scene is rendered as a Manim animation using the selected template. "
        "<b>✏️ Edit</b> to change the voiceover, scene type, or duration. "
        "<b>🔄 Regen</b> to let AI redesign a scene. "
        "Scene order is fixed — click <b>🚀 Render</b> when ready."
        "</div>",
        unsafe_allow_html=True,
    )

    # Scene type breakdown
    type_counts = {}
    for i in range(n):
        t = _get_topic_scene(i).get("scene_type","?")
        type_counts[t] = type_counts.get(t, 0) + 1
    type_str = "  ·  ".join(f"`{k}` ×{v}" for k, v in sorted(type_counts.items()))
    st.markdown(f'<div style="font-size:11px;color:#475569;margin-bottom:16px">Scene types: {type_str}</div>', unsafe_allow_html=True)

    for i in range(n):
        _render_topic_scene_card(i)

    st.markdown('<div class="section-sep"></div>', unsafe_allow_html=True)
    fin_bot = st.button("🚀 Render Topic Video", type="primary", use_container_width=True, key="topic_fin_bot")
    if fin_top or fin_bot:
        st.session_state.phase = "topic_rendering"; st.rerun()


# ════════════════════════════════════════════════════════════════════════════
# PHASE 3A — FULL PIPELINE RENDERING
# ════════════════════════════════════════════════════════════════════════════

def _phase_rendering():
    with st.sidebar:
        _show_logo_sidebar()
        st.markdown("## 🎓 Lecture Generator")
        st.markdown("---")
        st.info("⏳ Rendering in progress…\nPlease do not close this tab.")

    _show_logo_topright()
    st.title("🎬 Rendering Lecture Video")
    st.markdown('<div class="section-sep"></div>', unsafe_allow_html=True)

    progress   = st.progress(0, text="Preparing…")
    status_box = st.empty()
    log_ph     = st.empty()
    log_lines  = []
    handler    = _LogHandler(log_lines)
    handler.setFormatter(logging.Formatter("%(asctime)s  %(levelname)-7s  %(message)s"))
    logging.getLogger().addHandler(handler)

    def _refresh():
        html = "<br>".join(log_lines[-60:])
        log_ph.markdown(f'<div class="log-box">{html}</div>', unsafe_allow_html=True)

    def _upd(pct, msg):
        progress.progress(pct, text=msg); status_box.markdown(f"**{msg}**"); _refresh()

    try:
        _upd(5, "Assembling final slide list…")
        final_slides = _get_final_slides()
        n_ed   = sum(1 for v in st.session_state.slide_statuses.values() if v=="edited")
        n_reg  = sum(1 for v in st.session_state.slide_statuses.values() if v=="regenerated")
        n_cust = len(st.session_state.slide_custom_images)
        st.info(f"Rendering {len(final_slides)} slides ({n_ed} edited · {n_reg} AI-regen · {n_cust} custom images).")
        _upd(10, "Starting render pipeline…")
        start = time.time()
        out = create_lecture(
            topic=st.session_state.topic,
            transcript_text=st.session_state.transcript_text,
            slides_text=st.session_state.slides_text,
            notes_text=st.session_state.notes_text,
            skip_diagrams=st.session_state.skip_diagrams,
            slides_override=final_slides,
            custom_images_bytes=dict(st.session_state.slide_custom_images),
        )
        elapsed = time.time() - start
        progress.progress(100, text="✅ Complete!"); _refresh()
        st.session_state.output_path    = out
        st.session_state.render_elapsed = elapsed
        st.session_state.review_phase   = "review"
        st.session_state.phase          = "done"
        st.rerun()
    except Exception as e:
        progress.empty(); _refresh()
        st.error(f"❌ Rendering failed: {e}")
        with st.expander("🔍 Full traceback"): st.code(traceback.format_exc())
        col1, col2 = st.columns(2)
        with col1:
            if st.button("← Back to Review"): st.session_state.phase="review"; st.rerun()
        with col2:
            if st.button("🔄 Start Over"): st.session_state.phase="upload"; st.rerun()
    finally:
        logging.getLogger().removeHandler(handler)


# ════════════════════════════════════════════════════════════════════════════
# PHASE 3B — SLIDES PIPELINE RENDERING
# ════════════════════════════════════════════════════════════════════════════

def _phase_slides_rendering():
    with st.sidebar:
        _show_logo_sidebar()
        st.markdown("## 🎓 Lecture Generator")
        st.markdown("---")
        st.info("⏳ Rendering in progress…\nPlease do not close this tab.")

    _show_logo_topright()
    st.title("🎬 Rendering Slides + Voiceover Video")
    st.markdown('<div class="section-sep"></div>', unsafe_allow_html=True)

    progress   = st.progress(0, text="Preparing…")
    status_box = st.empty()
    log_ph     = st.empty()
    log_lines  = []
    handler    = _LogHandler(log_lines)
    handler.setFormatter(logging.Formatter("%(asctime)s  %(levelname)-7s  %(message)s"))
    logging.getLogger().addHandler(handler)

    def _refresh():
        html = "<br>".join(log_lines[-60:])
        log_ph.markdown(f'<div class="log-box">{html}</div>', unsafe_allow_html=True)

    def _upd(pct, msg):
        progress.progress(pct, text=msg); status_box.markdown(f"**{msg}**"); _refresh()

    try:
        _upd(5, "Assembling voiceover scripts…")
        final_slides = _get_final_slides_for_slides_pipeline()
        pptx_path    = st.session_state.get("slides_pptx_path")
        if not pptx_path or not os.path.exists(pptx_path):
            pptx_bytes = st.session_state.get("slides_pptx_bytes")
            if pptx_bytes:
                os.makedirs(OUTPUT_DIR, exist_ok=True)
                pptx_path = os.path.join(OUTPUT_DIR, "uploaded_slides.pptx")
                with open(pptx_path, "wb") as f: f.write(pptx_bytes)
                st.session_state.slides_pptx_path = pptx_path
            else:
                st.error("❌ PPTX file not found. Please start over and re-upload.")
                if st.button("🔄 Start Over"): st.session_state.phase="upload"; st.rerun()
                return
        n       = len(final_slides)
        est_min = max(1, sum(s.get("slide_duration_seconds",8) for s in final_slides)//60)
        st.info(f"Rendering {n} slides · ~{est_min} min · PPTX: {st.session_state.slides_pptx_filename or os.path.basename(pptx_path)}")
        _upd(10, "Starting render pipeline…")
        start = time.time()
        out = create_slides_lecture(
            pptx_path=pptx_path,
            transcript_text=st.session_state.get("transcript_text"),
            notes_text=st.session_state.get("notes_text"),
            context=st.session_state.get("slides_context",""),
            slides_override=final_slides,
        )
        elapsed = time.time() - start
        progress.progress(100, text="✅ Complete!"); _refresh()
        st.session_state.output_path    = out
        st.session_state.render_elapsed = elapsed
        st.session_state.review_phase   = "slides_review"
        st.session_state.phase          = "done"
        st.rerun()
    except Exception as e:
        progress.empty(); _refresh()
        st.error(f"❌ Rendering failed: {e}")
        with st.expander("🔍 Full traceback"): st.code(traceback.format_exc())
        col1, col2 = st.columns(2)
        with col1:
            if st.button("← Back to Review"): st.session_state.phase="slides_review"; st.rerun()
        with col2:
            if st.button("🔄 Start Over"): st.session_state.phase="upload"; st.rerun()
    finally:
        logging.getLogger().removeHandler(handler)


# ════════════════════════════════════════════════════════════════════════════
# PHASE 3C — TOPIC VIDEO RENDERING
# ════════════════════════════════════════════════════════════════════════════

def _phase_topic_rendering():
    with st.sidebar:
        _show_logo_sidebar()
        st.markdown("## 🎓 Lecture Generator")
        st.markdown("---")
        st.info("⏳ Rendering Manim scenes…\nThis may take several minutes.\nPlease do not close this tab.")

    _show_logo_topright()
    st.title("🎨 Rendering Topic Video (Pure Manim)")
    st.markdown('<div class="section-sep"></div>', unsafe_allow_html=True)

    progress   = st.progress(0, text="Preparing…")
    status_box = st.empty()
    log_ph     = st.empty()
    log_lines  = []
    handler    = _LogHandler(log_lines)
    handler.setFormatter(logging.Formatter("%(asctime)s  %(levelname)-7s  %(message)s"))
    logging.getLogger().addHandler(handler)

    def _refresh():
        html = "<br>".join(log_lines[-60:])
        log_ph.markdown(f'<div class="log-box">{html}</div>', unsafe_allow_html=True)

    def _upd(pct, msg):
        progress.progress(pct, text=msg); status_box.markdown(f"**{msg}**"); _refresh()

    try:
        _upd(5, "Assembling final scene list…")
        final_scenes = _get_final_topic_scenes()
        topic        = st.session_state.topic_topic
        context      = st.session_state.topic_context
        n            = len(final_scenes)
        total_dur    = sum(s.get("duration_seconds",30) for s in final_scenes)
        est_min      = max(1, total_dur // 60)

        n_ed  = sum(1 for v in st.session_state.topic_scene_statuses.values() if v=="edited")
        n_reg = sum(1 for v in st.session_state.topic_scene_statuses.values() if v=="regenerated")
        st.info(f"Rendering {n} Manim scenes (~{est_min} min) · {n_ed} edited · {n_reg} AI-regenerated · Topic: '{topic}'")

        _upd(10, "Starting Manim render + TTS generation in parallel…")
        start = time.time()

        out = create_topic_video(
            topic=topic,
            context=context,
            target_minutes=st.session_state.topic_minutes,
            scenes_override=final_scenes,
        )

        elapsed = time.time() - start
        progress.progress(100, text="✅ Complete!"); _refresh()
        st.session_state.output_path    = out
        st.session_state.render_elapsed = elapsed
        st.session_state.review_phase   = "topic_review"
        st.session_state.phase          = "done"
        st.rerun()

    except Exception as e:
        progress.empty(); _refresh()
        st.error(f"❌ Rendering failed: {e}")
        with st.expander("🔍 Full traceback"): st.code(traceback.format_exc())
        col1, col2 = st.columns(2)
        with col1:
            if st.button("← Back to Review"): st.session_state.phase="topic_review"; st.rerun()
        with col2:
            if st.button("🔄 Start Over"): st.session_state.phase="upload"; st.rerun()
    finally:
        logging.getLogger().removeHandler(handler)


# ════════════════════════════════════════════════════════════════════════════
# PHASE 4 — DONE  (shared by all pipelines)
# ════════════════════════════════════════════════════════════════════════════

def _phase_done():
    review_phase = st.session_state.get("review_phase","review")
    pipeline     = st.session_state.get("pipeline","full")

    with st.sidebar:
        _show_logo_sidebar()
        st.markdown("## 🎓 Lecture Generator")
        st.markdown("---")
        st.success("✅ Video ready!")
        st.markdown("<br>", unsafe_allow_html=True)
        if st.button("← Back to Review", use_container_width=True, key="done_back"):
            st.session_state.phase       = review_phase
            st.session_state.output_path = st.session_state.render_elapsed = None
            st.rerun()
        st.markdown("<br>", unsafe_allow_html=True)
        if st.button("🔄 Create New Lecture", use_container_width=True, key="done_new"):
            for k in list(st.session_state.keys()): del st.session_state[k]
            st.rerun()

    _show_logo_topright()
    st.title("🎉 Video Ready!")
    st.markdown('<div class="section-sep"></div>', unsafe_allow_html=True)

    elapsed   = st.session_state.render_elapsed or 0
    meta_path = os.path.join(OUTPUT_DIR, "lecture_metadata.json")
    meta      = {}
    if os.path.exists(meta_path):
        with open(meta_path) as f:
            meta = json.load(f).get("metadata", {})

    c1, c2, c3, c4 = st.columns(4)
    for col, val, lbl in [
        (c1, str(meta.get("total_slides","—")), "Scenes / Slides"),
        (c2, f"{meta.get('estimated_duration_minutes','—')} min", "Est. Duration"),
        (c3, str(meta.get("target_audience","—")).title(), "Audience"),
        (c4, f"{elapsed/60:.1f} min", "Render Time"),
    ]:
        with col:
            st.markdown(f'<div class="metric-card"><div class="metric-value">{val}</div><div class="metric-label">{lbl}</div></div>', unsafe_allow_html=True)

    if meta.get("full_description"):
        with st.expander("📝 Video Description"): st.write(meta["full_description"])

    # Pipeline-specific summary
    if pipeline == "full":
        n_ed   = sum(1 for v in st.session_state.slide_statuses.values() if v=="edited")
        n_reg  = sum(1 for v in st.session_state.slide_statuses.values() if v=="regenerated")
        n_cust = len(st.session_state.slide_custom_images)
        if n_ed + n_reg + n_cust > 0:
            st.info(f"**{n_ed}** manually edited · **{n_reg}** AI-regenerated · **{n_cust}** custom images")
    elif pipeline == "slides":
        fname_pptx = st.session_state.get("slides_pptx_filename","")
        n = len(st.session_state.get("slides_structure") or [])
        if fname_pptx: st.info(f"📎 **{fname_pptx}** · {n} slides · voiceover narration added")
    elif pipeline == "topic":
        n_ed  = sum(1 for v in st.session_state.topic_scene_statuses.values() if v=="edited")
        n_reg = sum(1 for v in st.session_state.topic_scene_statuses.values() if v=="regenerated")
        topic = st.session_state.get("topic_topic","")
        n     = len(st.session_state.get("topic_scenes") or [])
        st.info(f"⚡ **{topic}** · {n} pure Manim scenes · **{n_ed}** edited · **{n_reg}** AI-regen")

    st.markdown('<div class="section-sep"></div>', unsafe_allow_html=True)
    st.subheader("📥 Downloads")

    # Build download filename
    if pipeline == "full":
        fname = (st.session_state.topic[:35] or "lecture").replace(" ","_")
    elif pipeline == "slides":
        base  = st.session_state.get("slides_pptx_filename","lecture")
        fname = Path(base).stem.replace(" ","_")[:35]
    else:
        fname = (st.session_state.get("topic_topic","topic")[:35]).replace(" ","_")

    dc1, dc2, dc3 = st.columns(3)
    with dc1:
        out = st.session_state.output_path
        if out and os.path.exists(out):
            with open(out,"rb") as f:
                st.download_button("⬇️ Video (MP4)", data=f, file_name=f"{fname}_video.mp4",
                                   mime="video/mp4", use_container_width=True)
    with dc2:
        if os.path.exists(meta_path):
            with open(meta_path,"rb") as f:
                st.download_button("⬇️ Metadata (JSON)", data=f, file_name="lecture_metadata.json",
                                   mime="application/json", use_container_width=True)
    with dc3:
        desc_path = os.path.join(OUTPUT_DIR, "lecture_description.md")
        if os.path.exists(desc_path):
            with open(desc_path,"rb") as f:
                st.download_button("⬇️ Description (MD)", data=f, file_name="lecture_description.md",
                                   mime="text/markdown", use_container_width=True)

    st.markdown('<div class="section-sep"></div>', unsafe_allow_html=True)
    st.markdown(
        '<div class="phase-banner">'
        "💡 Want to tweak something? Click <b>← Back to Review</b> and re-render — it's fast because only changed scenes are re-processed."
        "</div>",
        unsafe_allow_html=True,
    )


# ════════════════════════════════════════════════════════════════════════════
# ROUTER
# ════════════════════════════════════════════════════════════════════════════

_PHASES = {
    "upload":           _phase_upload,
    "review":           _phase_review,
    "rendering":        _phase_rendering,
    "slides_review":    _phase_slides_review,
    "slides_rendering": _phase_slides_rendering,
    "topic_review":     _phase_topic_review,
    "topic_rendering":  _phase_topic_rendering,
    "done":             _phase_done,
}

_PHASES.get(st.session_state.get("phase", "upload"), _phase_upload)()
